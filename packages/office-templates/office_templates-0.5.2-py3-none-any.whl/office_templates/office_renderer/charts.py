from __future__ import annotations

from io import BytesIO
from itertools import zip_longest
from typing import TYPE_CHECKING, Callable, Optional

from pptx.chart.data import ChartData

from office_templates.templating.list import process_text_list

from .exceptions import ChartError
from .utils import get_load_workbook

if TYPE_CHECKING:
    from pptx.chart.chart import Chart


def process_chart(
    chart: Chart,
    context: dict,
    check_permissions: Optional[Callable[[object], bool]],
):
    """
    Read the current chart data, format text fields using the provided context,
    then update the chart with the new data.

    Parameters:
        chart: a pptx.chart.chart.Chart object.
        context: a dict used for formatting placeholders in text values.
        check_permissions: function to check permissions for objects.
    """

    # Common kwargs
    process_kwargs = dict(
        context=context,
        check_permissions=check_permissions,
        fail_if_empty=True,
    )

    # (1) Process series DATA in the chart (excluding series names and categories).
    all_processed_series_data = process_series_data(chart, context, check_permissions)
    # Transpose the data if the chart's axes are swapped.
    if chart_axes_are_swapped(chart):
        all_processed_series_data = list(zip(*all_processed_series_data))

    # (2) Process series names, and attach to the series data from before.
    raw_series_names = [series.name for series in chart.series]
    series_names = process_text_list(raw_series_names, **process_kwargs, as_float=False)
    series_names_values = [
        (series_name, series_data)
        for series_name, series_data in zip_longest(
            series_names,
            all_processed_series_data,
            fillvalue=None,
        )
    ]

    # (3) Process the placeholders in the categories from the first plot.
    if len(chart.plots) != 1:
        raise ChartError("Chart must have exactly one plot.")
    plot = chart.plots[0]
    raw_categories = [str(cat) for cat in plot.categories]
    categories = process_text_list(raw_categories, **process_kwargs, as_float=False)
    # Need at least one category, otherwise the chart will not display.
    if not categories:
        raise ChartError(
            f"Chart must have at least one category. Check this placeholder: {raw_categories}."
        )

    # (4) Create a new ChartData object and populate it.
    chart_data = ChartData()
    chart_data.categories = categories
    for name, values in series_names_values:
        chart_data.add_series(name, values)

    # (5) Replace the chart's data with the new ChartData.
    chart.replace_data(chart_data)


def get_raw_chart_data(chart: Chart):
    """
    Get the raw row/col data from a chart object.

    The return value is a list of lists, where outer list is COLUMNS and
    inner lists are ROWS.

    (This is because if we put text values into the series data, the pptx
    library will not be able to retrieve it (because it will be left out of
    the XML) and so we need to go to the underlying Excel workbook to update
    the data.)
    """

    load_workbook = get_load_workbook()

    # Access the embedded Excel workbook via the chart's part.
    chart_part = chart.part
    chart_workbook = chart_part.chart_workbook
    xlsx_part = chart_workbook.xlsx_part

    # Load the workbook from the blob.
    wb_stream = BytesIO(xlsx_part.blob)
    workbook = load_workbook(wb_stream)
    if workbook is None or not workbook.worksheets:
        raise ChartError("Attached Excel workbook is empty or has no worksheets.")
    worksheet = workbook.worksheets[0]

    # Retrieve the DATA from the workbook.
    return [[cell.value for cell in col] for col in worksheet.iter_cols()]


def process_series_data(
    chart: Chart,
    context: dict,
    check_permissions: Optional[Callable[[object], bool]],
):
    """
    Process the data in a chart, excluding the series names and categories.

    (This is because if we put text values into the series data, the pptx
    library will not be able to retrieve it (because it will be left out of
    the XML) and so we need to go to the underlying Excel workbook to update
    the data.)
    """

    # Get the raw data from the chart's workbook.
    # This is a list of lists, where the outer list is COLUMN and
    # the inner lists are items in each ROW.
    raw_data = get_raw_chart_data(chart)

    # Process the DATA in the workbook (excluding the the first row and column
    # which correspond to series/categories, which are handled more directly
    # instead of reading the Excel data).
    return [
        # Process placeholders in the list data.
        process_text_list(
            # Skip the first row, which is handled by the ChartData series/categories.
            col[1:],
            context=context,
            check_permissions=check_permissions,
            as_float=True,
            fail_if_not_float=True,
            fail_if_empty=False,
        )
        # Skip the first column, which is handled by the ChartData series/categories.
        for col in raw_data[1:]
    ]


def chart_axes_are_swapped(chart: Chart):
    """
    Return True if the X and Y axes are swapped in the chart.
    Check if the first series name is in the first column, which is the default.
    """

    # Look at the first series
    if not chart.series:
        return False
    series = chart.series[0]

    # Look at the cell range reference ("<c:f>") - but NOT for categories
    if series._element is None:
        return False
    range_references = series._element.xpath("./c:tx//c:f/text()")
    if not range_references:
        return False
    range_reference = range_references[0]

    # Normally, series names are in row 1, and if the row/column are
    # swapped, the series names are in column 1 (i.e. column "A"), so
    # if the range reference contains "$A$" then the axes ARE swapped.
    return "$A$" in range_reference

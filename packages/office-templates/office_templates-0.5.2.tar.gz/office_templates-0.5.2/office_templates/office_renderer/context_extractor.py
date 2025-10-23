import re
from typing import IO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .charts import get_raw_chart_data
from .constants import LOOP_START_PATTERN_STR
from .pptx.loops import extract_loop_directive
from .paragraphs import merge_split_placeholders
from .utils import get_load_workbook, identify_file_type

# Pattern to match placeholders, e.g. "{{ some.placeholder }}"
PLACEHOLDER_PATTERN = re.compile(r"\{\{\s*(.*?)\s*\}\}")

# Loop pattern is imported from constants
LOOP_START_PATTERN = re.compile(LOOP_START_PATTERN_STR)


def extract_top_level_context_keys_from_text(text: str) -> dict[str, list[str]]:
    """
    Given a text string, find all placeholders and extract the top-level context key.
    Returns a dict with:
        - simple_fields: keys without square brackets or periods in the placeholder
        - object_fields: keys where the placeholder includes a square bracket or period
    """
    KEYS_TO_IGNORE = {"now", "loop_count", "loop_number"}

    simple_fields = set()
    object_fields = set()
    placeholders = PLACEHOLDER_PATTERN.findall(text)
    for ph in placeholders:
        ph = ph.strip()
        if ph:
            m = re.match(r"([^\.\[\]\|]+)", ph)
            if m:
                key = m.group(1).strip()
                if key in KEYS_TO_IGNORE:
                    continue
                if ("." in ph) or ("[" in ph):
                    object_fields.add(key)
                else:
                    simple_fields.add(key)
    return {
        "simple_fields": sorted(simple_fields),
        "object_fields": sorted(object_fields),
    }


def _extract_texts_from_shape(shape, loop_variables: set) -> list[str]:
    """
    Recursively extract texts from a shape, handling grouped shapes.
    
    Args:
        shape: The shape to process
        loop_variables: Set to collect loop variables to ignore later
        
    Returns:
        List of text strings found in the shape
    """
    texts = []
    
    # Handle grouped shapes recursively
    if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for grouped_shape in shape.shapes:
            texts.extend(_extract_texts_from_shape(grouped_shape, loop_variables))
        return texts
    
    # Check for loop directives
    if hasattr(shape, "text_frame"):
        loop_var, loop_collection = extract_loop_directive(shape.text_frame.text)
        if loop_var and loop_collection:
            # Add loop variable to ignored list
            loop_variables.add(loop_var)

    # Process text frames.
    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            merge_split_placeholders(paragraph)
            texts.append(paragraph.text)
            
    # Process table cells.
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        merge_split_placeholders(paragraph)
                        texts.append(paragraph.text)
                        
    # Process chart spreadsheets.
    if getattr(shape, "has_chart", False):
        raw_data = get_raw_chart_data(shape.chart)
        for col in raw_data:
            for item in col:
                texts.append(str(item))
    
    return texts


def extract_context_keys_from_xlsx(template: str | IO[bytes]) -> dict[str, list[str]]:
    """
    Iterate through all worksheets and cells in the XLSX file (from a file path or file-like object),
    and return a dict with:
        - simple_fields: sorted list of unique simple keys
        - object_fields: sorted list of unique object keys
    """
    load_workbook = get_load_workbook()

    # Load the workbook (support template as a file path or file-like object)
    if isinstance(template, str):
        workbook = load_workbook(template)
    else:
        template.seek(0)
        workbook = load_workbook(template)

    simple_fields = set()
    object_fields = set()

    # Build a list of all cell values from all worksheets
    texts = []
    for sheet_name in workbook.sheetnames:
        worksheet = workbook[sheet_name]

        # Iterate through all cells that have values
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.value is not None:
                    texts.append(str(cell.value))

    # Process all texts to extract context keys
    for text in texts:
        keys = extract_top_level_context_keys_from_text(text)
        simple_fields.update(keys["simple_fields"])
        object_fields.update(keys["object_fields"])

    return {
        "simple_fields": sorted(simple_fields),
        "object_fields": sorted(object_fields),
    }


def extract_context_keys_from_pptx(template: str | IO[bytes]) -> dict[str, list[str]]:
    """
    Iterate through all slides, shapes, paragraphs and table cells in the PPTX (from a file path or file-like object),
    merging split placeholders, and return a dict with:
        - simple_fields: sorted list of unique simple keys
        - object_fields: sorted list of unique object keys
    """
    if isinstance(template, str):
        prs = Presentation(template)
    else:
        template.seek(0)
        prs = Presentation(template)

    simple_fields = set()
    object_fields = set()
    loop_variables = set()  # Store loop iterator variables to ignore them

    # Build a list of all texts on all slides.
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # Use the helper function to extract texts from shape (handles grouped shapes)
            shape_texts = _extract_texts_from_shape(shape, loop_variables)
            texts.extend(shape_texts)

    # Process all texts to extract context keys.
    for text in texts:
        keys = extract_top_level_context_keys_from_text(text)
        simple_fields.update(keys["simple_fields"])
        object_fields.update(keys["object_fields"])

    # Remove loop variables from the extracted fields
    object_fields = object_fields - loop_variables

    return {
        "simple_fields": sorted(simple_fields),
        "object_fields": sorted(object_fields),
    }


def extract_context_keys(template: str | IO[bytes]) -> dict[str, list[str]]:
    """
    Iterate through all slides/worksheets and their content in PPTX/XLSX files (from a file path or file-like object),
    and return a dict with:
        - simple_fields: sorted list of unique simple keys
        - object_fields: sorted list of unique object keys

    This function automatically detects the file type and uses the appropriate extraction method.
    """
    # Identify the file type
    file_type = identify_file_type(template)

    if file_type == "pptx":
        return extract_context_keys_from_pptx(template)
    elif file_type == "xlsx":
        return extract_context_keys_from_xlsx(template)
    else:
        # This should not happen due to identify_file_type raising UnsupportedFileType
        # but keeping for safety
        raise ValueError(f"Unsupported file type: {file_type}")

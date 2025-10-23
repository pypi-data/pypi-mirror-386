import unittest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from office_templates.office_renderer.tables import process_table_cell


class TestTablesIntegration(unittest.TestCase):
    def test_pure_placeholder_in_real_table(self):
        # 1) Create a presentation and add a slide
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # 2) Add a table with 1 row, 2 columns
        shapes = slide.shapes
        rows, cols = 1, 2
        left, top, width, height = 100000, 100000, 5000000, 200000
        table_shape = shapes.add_table(rows, cols, left, top, width, height)
        # 3) Grab the cell we want to test
        cell = table_shape.table.cell(0, 0)

        # 4) Set a placeholder text
        cell.text = "{{ test }}"

        # 5) Prepare context
        context = {"test": "HelloWorld"}

        # 6) Call process_table_cell with real objects
        process_table_cell(cell, context, None)

        # 7) Verify that placeholder was replaced
        self.assertEqual(cell.text, "HelloWorld")

        # Optional: check that the shape is truly a table
        self.assertEqual(table_shape.shape_type, MSO_SHAPE_TYPE.TABLE)


if __name__ == "__main__":
    unittest.main()

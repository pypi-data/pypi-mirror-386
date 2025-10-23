import os
import io
import tempfile
import unittest
from unittest.mock import patch, MagicMock

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import render_pptx


class TestPptxIntegrationLoops(unittest.TestCase):
    def setUp(self):
        # Create a temporary file for the PPTX
        self.temp_pptx = tempfile.mktemp(suffix=".pptx")
        self.output_pptx = tempfile.mktemp(suffix=".pptx")

        # Create a presentation with loop directive
        self.prs = Presentation()

        # Create slides for the loop
        self.slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.slide3 = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Add loop start shape to slide 1
        loop_start = self.slide1.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in users%"

        # Add content shape to slide 2
        user_info = self.slide2.shapes.add_textbox(
            Inches(1), Inches(2), Inches(4), Inches(1)
        )
        user_info.text_frame.text = "Name: {{ user.name }}, Email: {{ user.email }}"

        # Add loop end shape to slide 3
        loop_end = self.slide3.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Save the presentation
        self.prs.save(self.temp_pptx)

        # Prepare the context
        self.context = {
            "users": [
                {"name": "Alice", "email": "alice@example.com"},
                {"name": "Bob", "email": "bob@example.com"},
                {"name": "Charlie", "email": "charlie@example.com"},
            ]
        }

    def tearDown(self):
        # Clean up temporary files
        for temp_file in [self.temp_pptx, self.output_pptx]:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    @patch("office_templates.office_renderer.pptx.render.Presentation")
    def test_process_loops_integration(self, mock_presentation):
        """Test the integration of loop processing with minimal mocking."""
        from office_templates.office_renderer.pptx.loops import process_loops

        context = {"users": ["Alice", "Bob", "Charlie"]}
        errors = []

        prs = MagicMock()
        slides_list = []

        for i in range(3):
            slide = MagicMock()
            shapes = []

            if i == 0:  # First slide with loop start
                shape = MagicMock()
                shape.text_frame.text = "%loop user in users%"
                shapes.append(shape)
            elif i == 2:  # Last slide with loop end
                shape = MagicMock()
                shape.text_frame.text = "%endloop%"
                shapes.append(shape)
            else:  # Middle slide with content
                shape = MagicMock()
                shape.text_frame.text = "User: {{ user }}"
                shapes.append(shape)

            slide.shapes = shapes
            # Add slide_layout attribute for duplicate_slide
            slide.slide_layout = MagicMock()
            slides_list.append(slide)

        # Mock the slides collection with add_slide and __getitem__/__iter__
        slides_mock = MagicMock()
        slides_mock.__iter__.return_value = iter(slides_list)
        slides_mock.__getitem__.side_effect = slides_list.__getitem__
        slides_mock.__len__.return_value = len(slides_list)

        def add_slide(layout):
            # Return a new MagicMock slide with the given layout
            new_slide = MagicMock()
            new_slide.slide_layout = layout
            # Mock shapes with _spTree attribute
            shapes_mock = MagicMock()
            shapes_mock._spTree = MagicMock()
            new_slide.shapes = shapes_mock
            return new_slide

        slides_mock.add_slide.side_effect = add_slide
        prs.slides = slides_mock

        # Add a patch to detect loop directives during the test
        with patch(
            "office_templates.office_renderer.pptx.loops.is_loop_start"
        ) as mock_is_start:
            with patch(
                "office_templates.office_renderer.pptx.loops.is_loop_end"
            ) as mock_is_end:
                # Configure mocks to return true only for specific slides
                mock_is_start.side_effect = (
                    lambda shape: shape.text_frame.text.startswith("%loop")
                )
                mock_is_end.side_effect = lambda shape: shape.text_frame.text.startswith(
                    "%endloop"
                )

                # Process loops
                result = process_loops(prs, context, None, errors)

                # Check that we get 9 slides in the result (3 slides × 3 users)
                self.assertEqual(len(result), 9)

                # Check that there are no errors
                self.assertEqual(len(errors), 0)

    @patch("office_templates.office_renderer.pptx.render.Presentation")
    def test_dot_notation_processing(self, mock_presentation):
        """Test dot notation processing with minimal mocking."""
        # Setup a simplified test focusing on dot notation
        from office_templates.office_renderer.pptx.loops import process_loops, resolve_tag

        # Prepare context with dot notation
        context = {"program": {"members": ["Alice", "Bob"]}}
        errors = []

        # Create a minimalist presentation with 3 slides (start, content, end)
        prs = MagicMock()
        slides_list = []

        # Create slides with appropriate shapes
        for i in range(3):
            slide = MagicMock()
            shapes = []

            if i == 0:  # First slide with loop start using dot notation
                shape = MagicMock()
                shape.text_frame.text = "%loop member in program.members%"
                shapes.append(shape)
            elif i == 2:  # Last slide with loop end
                shape = MagicMock()
                shape.text_frame.text = "%endloop%"
                shapes.append(shape)
            else:  # Middle slide with content
                shape = MagicMock()
                shape.text_frame.text = "Member: {{ member }}"
                shapes.append(shape)

            slide.shapes = shapes
            # Add slide_layout attribute for duplicate_slide
            slide.slide_layout = MagicMock()
            slides_list.append(slide)

        # Mock the slides collection with add_slide and __getitem__/__iter__
        slides_mock = MagicMock()
        slides_mock.__iter__.return_value = iter(slides_list)
        slides_mock.__getitem__.side_effect = slides_list.__getitem__
        slides_mock.__len__.return_value = len(slides_list)

        def add_slide(layout):
            # Return a new MagicMock slide with the given layout
            new_slide = MagicMock()
            new_slide.slide_layout = layout
            # Mock shapes with _spTree attribute
            shapes_mock = MagicMock()
            shapes_mock._spTree = MagicMock()
            new_slide.shapes = shapes_mock
            return new_slide

        slides_mock.add_slide.side_effect = add_slide
        prs.slides = slides_mock

        # Add a patch to detect loop directives during the test
        with patch(
            "office_templates.office_renderer.pptx.loops.is_loop_start"
        ) as mock_is_start:
            with patch(
                "office_templates.office_renderer.pptx.loops.is_loop_end"
            ) as mock_is_end:
                # Configure mocks to return true only for specific slides
                mock_is_start.side_effect = (
                    lambda shape: shape.text_frame.text.startswith("%loop")
                )
                mock_is_end.side_effect = lambda shape: shape.text_frame.text.startswith(
                    "%endloop"
                )

                # Test direct resolution of the dot notation path
                members = resolve_tag("program.members", context, None)
                self.assertEqual(members, ["Alice", "Bob"])

                # Process loops
                result = process_loops(prs, context, None, errors)

                # Check that we get 6 slides in the result (3 slides × 2 members)
                self.assertEqual(len(result), 6)

                # Check that there are no errors
                self.assertEqual(len(errors), 0)


if __name__ == "__main__":
    unittest.main()

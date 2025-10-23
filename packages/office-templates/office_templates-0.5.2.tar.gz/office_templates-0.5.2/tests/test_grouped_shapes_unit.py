"""
Unit tests for grouped shapes functionality in PPTX rendering.
"""

import tempfile
import unittest
from unittest.mock import Mock

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from office_templates.office_renderer.pptx.render import process_shape_content


class TestGroupedShapesUnit(unittest.TestCase):
    def setUp(self):
        """Create a test presentation with grouped shapes."""
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Create some individual shapes to group
        self.text_box1 = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(1)
        )
        self.text_box1.text_frame.text = "{{ grouped_variable }}"
        
        self.text_box2 = self.slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(2), Inches(1)
        )
        self.text_box2.text_frame.text = "Static text: {{ another_var }}"
        
        # Group the shapes
        self.group = self.slide.shapes.add_group_shape([self.text_box1, self.text_box2])
        
        # Test context
        self.context = {
            "grouped_variable": "Rendered grouped text",
            "another_var": "Rendered another text",
        }
        
        # Mock objects for testing
        self.mock_check_permissions = Mock(return_value=True)
        self.errors = []

    def test_group_shape_detection(self):
        """Test that we can detect grouped shapes correctly."""
        # Find the group shape
        group_shape = None
        for shape in self.slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_shape = shape
                break
        
        self.assertIsNotNone(group_shape, "Should find a group shape")
        self.assertEqual(len(group_shape.shapes), 2, "Group should contain 2 shapes")
        
        # Verify the grouped shapes have the expected text
        texts = []
        for shape in group_shape.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                texts.append(shape.text_frame.text)
        
        self.assertIn("{{ grouped_variable }}", texts)
        self.assertIn("Static text: {{ another_var }}", texts)

    def test_process_grouped_shapes_now_works(self):
        """Test that the updated implementation processes grouped shapes correctly."""
        # Find the group shape
        group_shape = None
        for shape in self.slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_shape = shape
                break
        
        # Process the group shape with updated implementation
        process_shape_content(
            shape=group_shape,
            slide=self.slide,
            context=self.context,
            slide_number=1,
            check_permissions=self.mock_check_permissions,
            errors=self.errors,
        )
        
        # Check that variables inside the group were processed
        group_texts = []
        for shape in group_shape.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                group_texts.append(shape.text_frame.text)
        
        # Variables should now be processed
        self.assertIn("Rendered grouped text", group_texts)
        self.assertTrue(any("Rendered another text" in text for text in group_texts))
        
        # Original variables should be replaced
        self.assertFalse(any("{{ grouped_variable }}" in text for text in group_texts))
        self.assertFalse(any("{{ another_var }}" in text for text in group_texts))

    def test_nested_group_shapes(self):
        """Test handling of nested grouped shapes."""
        # Create another shape to add to a nested group
        text_box3 = self.slide.shapes.add_textbox(
            Inches(4), Inches(1), Inches(2), Inches(1)
        )
        text_box3.text_frame.text = "{{ nested_var }}"
        
        # Create a nested group (group within a group)
        # Note: This might not be directly possible via python-pptx API
        # but we should handle it if it exists in loaded files
        
        # For now, just test that we handle the concept
        self.context["nested_var"] = "Nested rendered text"
        
        # This test will be expanded once we implement recursive processing

    def tearDown(self):
        """Clean up any temporary files."""
        # No temp files in unit tests, but good practice
        pass


if __name__ == "__main__":
    unittest.main()
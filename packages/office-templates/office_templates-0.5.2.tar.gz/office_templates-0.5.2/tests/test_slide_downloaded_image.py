import os
import tempfile
import unittest
from unittest.mock import patch

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from PIL import Image

from office_templates.office_renderer.images import replace_shape_with_image, ImageError


class TestSlideDownloadedImage(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.textbox = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(2)
        )
        # Create a small red image on disk
        img = Image.new("RGB", (4, 4), color="red")
        self.temp_image = tempfile.mktemp(suffix=".png")
        img.save(self.temp_image)
        self.textbox.text_frame.text = f"%imagesqueeze% file://{self.temp_image}"
        self.context = {"img": f"file://{self.temp_image}"}

    def tearDown(self):
        if os.path.exists(self.temp_image):
            os.remove(self.temp_image)

    def test_replace_shape_with_image(self):
        original_width = self.textbox.width
        original_height = self.textbox.height

        replace_shape_with_image(self.textbox, self.slide, context=self.context)

        # Only one non-placeholder shape should remain and it should be a picture
        shapes = [
            shape for shape in self.slide.shapes if not shape.is_placeholder
        ]
        self.assertEqual(len(shapes), 1)
        pic = shapes[0]
        self.assertEqual(pic.shape_type, MSO_SHAPE_TYPE.PICTURE)
        self.assertEqual(pic.width, original_width)
        self.assertEqual(pic.height, original_height)

    @patch("office_templates.office_renderer.images.urlopen")
    def test_invalid_url_raises_image_error(self, mock_urlopen):
        mock_urlopen.side_effect = Exception("boom")
        self.textbox.text_frame.text = "%image% http://example.com/foo.png"
        with self.assertRaises(ImageError):
            replace_shape_with_image(self.textbox, self.slide, context={})

    def test_image_aspect_ratio_fitting(self):
        """%image% should keep aspect ratio inside the shape."""

        # Create non-square shape
        shape = self.slide.shapes.add_textbox(Inches(3), Inches(3), Inches(4), Inches(1))
        # Use explicit braces so the URL placeholder is passed through for processing
        shape.text_frame.text = "%image% {{ img }}"

        replace_shape_with_image(shape, self.slide, context=self.context)

        pic = [s for s in self.slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE][-1]
        self.assertEqual(pic.height, shape.height)
        self.assertLess(pic.width, shape.width)

    def test_placeholder_in_url(self):
        """URL expressions should be processed using process_text."""

        shape = self.slide.shapes.add_textbox(Inches(5), Inches(1), Inches(2), Inches(2))
        shape.text_frame.text = "%imagesqueeze% {{ img }}"

        replace_shape_with_image(shape, self.slide, context=self.context)

        pic = [s for s in self.slide.shapes if s.left == shape.left and s.top == shape.top and s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        self.assertEqual(len(pic), 1)


if __name__ == "__main__":
    unittest.main()
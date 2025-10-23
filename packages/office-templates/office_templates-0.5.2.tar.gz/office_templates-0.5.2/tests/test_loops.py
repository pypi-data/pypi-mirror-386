import unittest
from unittest.mock import MagicMock, patch

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from office_templates.office_renderer.pptx.loops import (
    extract_loop_directive,
    is_loop_end,
    is_loop_start,
    process_loops,
    LOOP_START_PATTERN,
    LOOP_END_PATTERN,
)


class TestLoopDirectives(unittest.TestCase):
    def setUp(self):
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.textbox = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(2)
        )

    def test_extract_loop_directive(self):
        """Test the extraction of loop variable and collection from directive."""
        # Valid loop directive
        self.textbox.text_frame.text = "%loop user in users%"
        variable, collection = extract_loop_directive(self.textbox.text_frame.text)
        self.assertEqual(variable, "user")
        self.assertEqual(collection, "users")

        # Valid loop directive with spaces
        self.textbox.text_frame.text = "% loop user in users %"
        variable, collection = extract_loop_directive(self.textbox.text_frame.text)
        self.assertEqual(variable, "user")
        self.assertEqual(collection, "users")

        # Valid loop directive with dot notation
        self.textbox.text_frame.text = "%loop user in program.users%"
        variable, collection = extract_loop_directive(self.textbox.text_frame.text)
        self.assertEqual(variable, "user")
        self.assertEqual(collection, "program.users")

        # Valid loop directive with dot notation and filtering
        self.textbox.text_frame.text = "%loop user in program.users[is_active=True]%"
        variable, collection = extract_loop_directive(self.textbox.text_frame.text)
        self.assertEqual(variable, "user")
        self.assertEqual(collection, "program.users[is_active=True]")

        # Invalid loop directive
        self.textbox.text_frame.text = "Not a loop directive"
        variable, collection = extract_loop_directive(self.textbox.text_frame.text)
        self.assertIsNone(variable)
        self.assertIsNone(collection)

        # Empty text
        variable, collection = extract_loop_directive(None)
        self.assertIsNone(variable)
        self.assertIsNone(collection)

    def test_is_loop_start(self):
        """Test the detection of loop start directive."""
        # Valid loop start
        self.textbox.text_frame.text = "%loop user in users%"
        self.assertTrue(is_loop_start(self.textbox))

        # Not a loop start
        self.textbox.text_frame.text = "Not a loop start"
        self.assertFalse(is_loop_start(self.textbox))

        # Shape has no text_frame - use proper mocking
        shape_without_text_frame = MagicMock(spec=[])  # No text_frame attribute
        self.assertFalse(is_loop_start(shape_without_text_frame))

        # Text frame has no text attribute - use proper mocking
        shape_with_empty_text_frame = MagicMock()
        text_frame_without_text = MagicMock(spec=[])  # No text attribute
        shape_with_empty_text_frame.text_frame = text_frame_without_text
        self.assertFalse(is_loop_start(shape_with_empty_text_frame))

    def test_is_loop_end(self):
        """Test the detection of loop end directive."""
        # Valid loop end
        self.textbox.text_frame.text = "%endloop%"
        self.assertTrue(is_loop_end(self.textbox))

        # Valid loop end with spaces
        self.textbox.text_frame.text = "% endloop %"
        self.assertTrue(is_loop_end(self.textbox))

        # Not a loop end
        self.textbox.text_frame.text = "Not a loop end"
        self.assertFalse(is_loop_end(self.textbox))

        # Shape has no text_frame - use proper mocking
        shape_without_text_frame = MagicMock(spec=[])  # No text_frame attribute
        self.assertFalse(is_loop_end(shape_without_text_frame))

        # Text frame has no text attribute - use proper mocking
        shape_with_empty_text_frame = MagicMock()
        text_frame_without_text = MagicMock(spec=[])  # No text attribute
        shape_with_empty_text_frame.text_frame = text_frame_without_text
        self.assertFalse(is_loop_end(shape_with_empty_text_frame))

    def test_regex_patterns(self):
        """Test that the regex patterns handle spaces correctly."""
        # Test LOOP_START_PATTERN
        self.assertIsNotNone(LOOP_START_PATTERN.search("%loop user in users%"))
        self.assertIsNotNone(LOOP_START_PATTERN.search("% loop user in users %"))
        self.assertIsNotNone(LOOP_START_PATTERN.search("%  loop  user  in  users  %"))

        # Test LOOP_END_PATTERN
        self.assertIsNotNone(LOOP_END_PATTERN.search("%endloop%"))
        self.assertIsNotNone(LOOP_END_PATTERN.search("% endloop %"))
        self.assertIsNotNone(LOOP_END_PATTERN.search("%  endloop  %"))


class TestLoopProcessing(unittest.TestCase):
    def setUp(self):
        self.context = {
            "users": [
                {"name": "Alice", "email": "alice@example.com"},
                {"name": "Bob", "email": "bob@example.com"},
            ]
        }

        # Create more complex context for testing dot notation
        self.complex_context = {
            "program": {
                "users": [
                    {"name": "Alice", "email": "alice@example.com", "is_active": True},
                    {"name": "Bob", "email": "bob@example.com", "is_active": True},
                    {
                        "name": "Charlie",
                        "email": "charlie@example.com",
                        "is_active": False,
                    },
                ]
            }
        }

        # Create a presentation and errors list for testing
        self.prs = Presentation()
        self.errors = []

    def test_process_loops_slide_numbers(self):
        """Test that process_loops correctly numbers slides."""
        # Create a presentation with 5 slides:
        # Slide 1: Normal slide
        # Slide 2: Loop start
        # Slide 3: Content inside loop
        # Slide 4: Loop end
        # Slide 5: Normal slide

        # Add all slides
        slides = []
        for i in range(5):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add loop start to slide 2
        loop_start = slides[1].shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in users%"

        # Add loop end to slide 4
        loop_end = slides[3].shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Process loops
        result = process_loops(self.prs, self.context, None, self.errors)

        # We should have:
        # - Slide 1 (normal) with slide_number 1
        # - Slides 2-4 (loop) repeated twice (2 users), with slide_numbers 2-7
        # - Slide 5 (normal) with slide_number 8

        # Check total number of slides to process
        self.assertEqual(len(result), 8)

        # Check slide numbers
        self.assertEqual(result[0]["slide_number"], 1)  # First normal slide

        # First loop iteration (3 slides)
        self.assertEqual(result[1]["slide_number"], 2)
        self.assertEqual(result[2]["slide_number"], 3)
        self.assertEqual(result[3]["slide_number"], 4)

        # Second loop iteration (3 slides)
        self.assertEqual(result[4]["slide_number"], 5)
        self.assertEqual(result[5]["slide_number"], 6)
        self.assertEqual(result[6]["slide_number"], 7)

        # Last normal slide
        self.assertEqual(result[7]["slide_number"], 8)

    def test_process_loops_dot_notation(self):
        """Test process_loops with dot notation for the collection."""
        # Create a presentation with 3 slides (start, content, end)
        slides = []
        for i in range(3):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add loop with dot notation to first slide
        loop_start = slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in program.users%"

        # Add loop end to last slide
        loop_end = slides[2].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Process loops
        result = process_loops(self.prs, self.complex_context, None, self.errors)

        # We should have 9 slides (3 slides × 3 users)
        self.assertEqual(len(result), 9)
        self.assertEqual(len(self.errors), 0)

    def test_process_loops_with_filtering(self):
        """Test process_loops with filtering in the collection."""
        # Create a presentation with 3 slides (start, content, end)
        slides = []
        for i in range(3):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add loop with filtering to first slide
        loop_start = slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in program.users[is_active=True]%"

        # Add loop end to last slide
        loop_end = slides[2].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Process loops
        result = process_loops(self.prs, self.complex_context, None, self.errors)

        # We should have 6 slides (3 slides × 2 active users)
        self.assertEqual(len(result), 6)
        self.assertEqual(len(self.errors), 0)

    def test_process_loops_with_nonexistent_collection(self):
        """Test process_loops with a nonexistent collection."""
        # Create a presentation with 3 slides (start, content, end)
        slides = []
        for i in range(3):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add loop with nonexistent collection
        loop_start = slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in nonexistent%"

        # Add loop end to last slide
        loop_end = slides[2].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Process loops
        result = process_loops(self.prs, self.complex_context, None, self.errors)

        # We should have an error but non-loop slides should be processed
        self.assertTrue(any("nonexistent" in error for error in self.errors))

    def test_process_loops_with_multiple_directives_on_same_slide(self):
        """Test process_loops with multiple loop directives on the same slide."""
        # Create a presentation with 3 slides
        slides = []
        for i in range(3):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add multiple loop starts to first slide
        loop_start1 = slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_start1.text_frame.text = "%loop user in users%"

        loop_start2 = slides[0].shapes.add_textbox(
            Inches(1), Inches(2), Inches(3), Inches(0.5)
        )
        loop_start2.text_frame.text = "%loop item in items%"

        # Add loop end to last slide
        loop_end = slides[2].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Process loops
        result = process_loops(self.prs, self.context, None, self.errors)

        # We should get errors about multiple loop starts
        self.assertTrue(
            any(
                "Multiple loop start directives on same slide" in error
                for error in self.errors
            )
        )

    def test_process_loops_with_both_start_and_end_on_same_slide(self):
        """Test process_loops with both start and end directives on the same slide."""
        # Create a presentation with 3 slides
        slides = []
        for i in range(3):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add both start and end to first slide
        loop_start = slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in users%"

        loop_end = slides[0].shapes.add_textbox(
            Inches(1), Inches(2), Inches(3), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Process loops
        result = process_loops(self.prs, self.context, None, self.errors)

        # We should NOT get an error about both directives on same slide
        self.assertEqual(len(self.errors), 0)

    def test_process_loops_with_unclosed_loop(self):
        """Test process_loops with an unclosed loop."""
        # Create a presentation with 3 slides
        slides = []
        for i in range(3):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add loop start but no end
        loop_start = slides[0].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in users%"

        # Process loops
        result = process_loops(self.prs, self.context, None, self.errors)

        # We should get an error about unclosed loop
        self.assertTrue(
            any("Loop started but never closed" in error for error in self.errors)
        )

    def test_process_loops_with_end_without_start(self):
        """Test process_loops with an end directive without a matching start."""
        # Create a presentation with 3 slides
        slides = []
        for i in range(3):
            slides.append(self.prs.slides.add_slide(self.prs.slide_layouts[5]))

        # Add loop end without start
        loop_end = slides[1].shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Process loops
        result = process_loops(self.prs, self.context, None, self.errors)

        # We should get an error about endloop without matching start
        self.assertTrue(
            any("without a matching loop start" in error for error in self.errors)
        )

    @patch("office_templates.office_renderer.pptx.render.process_loops")
    @patch("office_templates.office_renderer.pptx.render.Presentation")
    def test_process_loops_called(self, mock_presentation, mock_process_loops):
        """Test that process_loops is called from render_pptx."""
        from office_templates.office_renderer.pptx.render import render_pptx

        # Prepare mock return value for process_loops
        mock_process_loops.return_value = []

        # Setup mock for Presentation
        mock_prs = MagicMock()
        mock_presentation.return_value = mock_prs

        # Call render_pptx with a string template path
        render_pptx("template.pptx", {}, "output.pptx", None)

        # Verify process_loops was called
        mock_process_loops.assert_called_once()

    def test_extract_loop_directive_integration(self):
        """Test extract_loop_directive directly."""
        # Valid loop
        directive = "%loop user in users%"
        variable, collection = extract_loop_directive(directive)
        self.assertEqual(variable, "user")
        self.assertEqual(collection, "users")

        # Invalid format
        directive = "%loop invalid directive%"
        variable, collection = extract_loop_directive(directive)
        self.assertIsNone(variable)
        self.assertIsNone(collection)

        # Empty
        directive = ""
        variable, collection = extract_loop_directive(directive)
        self.assertIsNone(variable)
        self.assertIsNone(collection)


if __name__ == "__main__":
    unittest.main()

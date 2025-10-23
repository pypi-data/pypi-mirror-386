"""
Test the shape removal functionality.
"""

import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer.pptx.utils import remove_shape


class TestShapeRemoval(unittest.TestCase):
    def setUp(self):
        # Create a temporary file for the PPTX
        self.temp_pptx = tempfile.mktemp(suffix=".pptx")

        # Create a presentation with a textbox
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.textbox = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        self.textbox.text_frame.text = "Test text"

        # Save the presentation
        self.prs.save(self.temp_pptx)

    def tearDown(self):
        # Clean up temporary files
        if os.path.exists(self.temp_pptx):
            os.remove(self.temp_pptx)

    def test_remove_shape(self):
        """Test that remove_shape correctly removes a shape from a slide."""
        # Count initial shapes
        shapes = [shape for shape in self.slide.shapes if not shape.is_placeholder]
        initial_count = len(shapes)
        self.assertEqual(initial_count, 1, "Should start with 1 shape")

        # Remove the shape
        remove_shape(self.textbox)

        # Save and reload to ensure changes are applied
        self.prs.save(self.temp_pptx)
        prs = Presentation(self.temp_pptx)

        # Count shapes after removal
        shapes = [shape for shape in self.slide.shapes if not shape.is_placeholder]
        self.assertEqual(len(shapes), 0, "Shape should be removed")


if __name__ == "__main__":
    unittest.main()

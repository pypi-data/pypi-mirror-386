import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import compose_pptx


# Dummy objects for testing
class DummyUser:
    def __init__(self, name, email):
        self.name = name
        self.email = email
        self._meta = True

    def __str__(self):
        return self.name


class TestPPTXComposition(unittest.TestCase):
    def setUp(self):
        # Create test presentations with layouts and slides
        self.temp_files = []
        self.template1_path = self._create_test_template1()
        self.template2_path = self._create_test_template2()

        # Test context
        self.context = {
            "user": DummyUser("Alice", "alice@example.com"),
            "title": "Test Presentation",
            "content": "This is test content",
            "chart_title": "Sales Performance Chart",
            "product_summary": "Product performance overview",
        }

    def tearDown(self):
        # Clean up temporary files
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def _create_test_template1(self):
        """Create a test template with multiple layouts."""
        prs = Presentation()

        # Add a slide with title layout (for title layouts test)
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        slide1.shapes.title.text = "Title Layout"

        # Add a slide with content (for tagged layouts test)
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        slide2.shapes.title.text = "Content Slide"
        
        # Add %layout% tag in its own shape
        layout_box = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(0.5))
        layout_box.text_frame.text = "% layout content %"
        
        # Add template variables in separate shapes
        chart_title_box = slide2.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(0.5))
        chart_title_box.text_frame.text = "{{ chart_title }}"
        
        product_summary_box = slide2.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(0.5))
        product_summary_box.text_frame.text = "{{ product_summary }}"

        # Save to temp file
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def _create_test_template2(self):
        """Create another test template with different content."""
        prs = Presentation()

        # Add a slide for tagged layout
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        
        # Add %layout% tag in its own shape
        layout_box = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
        layout_box.text_frame.text = "% layout special %"
        
        # Add template variables in separate shapes
        title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(0.5))
        title_box.text_frame.text = "{{ title }}"
        
        content_box = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(0.5))
        content_box.text_frame.text = "{{ content }}"

        # Add another slide with title for title layouts
        slide2 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        slide2.shapes.title.text = "Special Layout"

        # Save to temp file
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def test_compose_with_master_layouts(self):
        """Test basic composition using master slide layouts."""
        slide_specs = [
            {"layout": "Title Slide", "title": "My Title"},
            {"layout": "Title and Content", "content": "Some content"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check output file exists and has correct number of slides
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_compose_with_tagged_layouts(self):
        """Test composition using tagged layouts."""
        slide_specs = [
            {"layout": "content", "title": "Content Title"},
            {"layout": "special", "content": "Special content"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path, self.template2_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

        # Verify that %layout% shapes are removed from output slides
        for slide in prs.slides:
            layout_shapes = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                    text = shape.text_frame.text.strip()
                    if "% layout" in text and "%" in text:
                        layout_shapes.append(shape)
            self.assertEqual(len(layout_shapes), 0, 
                           f"Found %layout% shapes in output slide: {[s.text_frame.text for s in layout_shapes]}")

    def test_compose_with_title_layouts(self):
        """Test composition using slide titles as layout IDs."""
        slide_specs = [
            {"layout": "Title Layout", "content": "From title layout"},
            {"layout": "Special Layout", "content": "From special layout"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path, self.template2_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_all_slides_as_layouts_by_title=True,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_compose_with_placeholders(self):
        """Test composition with placeholder processing."""
        slide_specs = [
            {
                "layout": "Title Slide",
                "placeholders": ["{{ title }}", "{{ user.name }}"],
                "title": "Placeholder Test",
            },
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_compose_with_missing_layout(self):
        """Test error handling when layout is not found."""
        slide_specs = [
            {"layout": "NonExistentLayout", "content": "Test"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should fail with errors
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(
            any("Layout 'NonExistentLayout' not found" in error for error in errors)
        )

    def test_compose_with_no_template_files_invalid_layout(self):
        """Test error handling when no template files and invalid layout is specified."""
        slide_specs = [{"layout": "test", "content": "Test"}]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should fail with errors because 'test' is not a valid default layout
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("Layout 'test' not found" in error for error in errors))

    def test_compose_with_no_slides(self):
        """Test error handling when no slides are specified."""
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slide_specs=[],
            global_context=self.context,
            output=output_file,
        )

        # Should fail with errors
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("No slides specified" in error for error in errors))

    def test_multiple_layout_shapes_error(self):
        """Test error handling when multiple %layout% shapes exist on same slide."""
        # Create a problematic template with multiple %layout% shapes
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        
        # Add two %layout% shapes with different IDs
        layout_box1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
        layout_box1.text_frame.text = "% layout test1 %"
        
        layout_box2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(0.5))
        layout_box2.text_frame.text = "% layout test2 %"
        
        # Save the problematic template
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        
        # Try to use this template - should fail
        slide_specs = [{"layout": "test1"}]
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[temp_file],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )
        
        # Should fail with errors
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("Multiple different layout IDs" in error for error in errors))

    def test_layout_with_loop_error(self):
        """Test error handling when %layout% and %loop% exist on same slide."""
        # Create a problematic template with both %layout% and %loop%
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        
        # Add %layout% shape
        layout_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
        layout_box.text_frame.text = "% layout test_layout %"
        
        # Add %loop% shape
        loop_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(0.5))
        loop_box.text_frame.text = "% loop item in items %"
        
        # Save the problematic template
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        
        # Try to use this template - should fail
        slide_specs = [{"layout": "test_layout"}]
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)
        
        result, errors = compose_pptx(
            template_files=[temp_file],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )
        
        # Should fail with errors
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("cannot contain %loop%" in error for error in errors))

    def test_compose_without_template_files(self):
        """Test composition without any template files using default layouts."""
        slide_specs = [
            {"layout": "Title Slide", "content": "Test slide with default layout"},
            {"layout": "Blank", "content": "Another test slide"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_compose_without_template_files_none(self):
        """Test composition with template_files=None (default parameter)."""
        slide_specs = [
            {"layout": "Title Slide", "content": "Test with None"},
            {"layout": "Title and Content", "content": "Second slide"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=None,
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_compose_without_template_files_with_placeholders(self):
        """Test composition without template files but with placeholder processing."""
        slide_specs = [
            {
                "layout": "Title Slide",
                "placeholders": ["{{ title }}", "{{ user.name }}"],
            },
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

    def test_compose_without_template_files_all_default_layouts(self):
        """Test that all default layouts are available when no template files provided."""
        # Use various default layouts
        slide_specs = [
            {"layout": "Title Slide"},
            {"layout": "Title and Content"},
            {"layout": "Section Header"},
            {"layout": "Two Content"},
            {"layout": "Comparison"},
            {"layout": "Title Only"},
            {"layout": "Blank"},
            {"layout": "Content with Caption"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 8)


if __name__ == "__main__":
    unittest.main()

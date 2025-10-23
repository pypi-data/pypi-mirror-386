import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import render_pptx
from office_templates.office_renderer.pptx.loops import is_loop_start, is_loop_end


class TestRealPptxLoops(unittest.TestCase):
    def setUp(self):
        """Create a real PPTX file with loop directives for testing."""
        # Create a temporary file for the PPTX
        self.temp_pptx = tempfile.mktemp(suffix=".pptx")
        self.output_pptx = tempfile.mktemp(suffix=".pptx")

        # Create a presentation with loop directive
        self.prs = Presentation()

        # Create slides for the loop
        self.slide1 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.slide2 = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self.slide3 = self.prs.slides.add_slide(self.prs.slide_layouts[5])

        # Add loop start shape to slide 1
        loop_start = self.slide1.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_start.text_frame.text = "%loop user in users%"

        # Add content shape to slide 2
        user_info = self.slide2.shapes.add_textbox(
            Inches(1), Inches(2), Inches(4), Inches(1)
        )
        user_info.text_frame.text = "Name: {{ user.name }}, Email: {{ user.email }}"

        # Add loop end shape to slide 3
        loop_end = self.slide3.shapes.add_textbox(
            Inches(1), Inches(1), Inches(2), Inches(0.5)
        )
        loop_end.text_frame.text = "%endloop%"

        # Save the presentation
        self.prs.save(self.temp_pptx)

        # Prepare the context
        self.context = {
            "users": [
                {"name": "Alice", "email": "alice@example.com"},
                {"name": "Bob", "email": "bob@example.com"},
                {"name": "Charlie", "email": "charlie@example.com"},
            ]
        }

    def tearDown(self):
        # Clean up temporary files
        for temp_file in [self.temp_pptx, self.output_pptx]:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def test_real_pptx_loop_functionality(self):
        """Test that loops work correctly in a real PPTX file."""
        # Render the template with the context
        output, errors = render_pptx(self.temp_pptx, self.context, self.output_pptx, None)

        # Ensure no errors occurred
        self.assertIsNone(errors, f"Errors occurred during rendering: {errors}")

        # Load the output presentation to verify results
        prs = Presentation(self.output_pptx)

        # Print all text content for debugging
        print("\nText content in all slides:")
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                    print(f"  Shape {j+1}: '{shape.text_frame.text}'")

        # Verify loop directive shapes have been deleted
        loop_directives_found = False
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                    text = shape.text_frame.text.strip()
                    # Check if any text looks like a loop directive
                    if (
                        text.startswith("%loop")
                        or text.startswith("% loop")
                        or text == "%endloop%"
                        or text == "% endloop %"
                    ):
                        loop_directives_found = True
                        break
            if loop_directives_found:
                break

        self.assertFalse(
            loop_directives_found, "Loop directives found with text not cleared"
        )

        # Verify at least one user's content was correctly inserted
        user_content_found = False
        for user in self.context["users"]:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                        if (
                            user["name"] in shape.text_frame.text
                            and user["email"] in shape.text_frame.text
                        ):
                            user_content_found = True
                            print(f"Found user content: {shape.text_frame.text}")

        self.assertTrue(user_content_found, "User content not found in any slide")


if __name__ == "__main__":
    unittest.main()

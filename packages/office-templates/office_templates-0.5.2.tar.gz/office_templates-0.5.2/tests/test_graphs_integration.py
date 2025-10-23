import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import compose_pptx


class TestGraphsIntegration(unittest.TestCase):
    """Integration tests for graph composition with realistic scenarios."""

    def setUp(self):
        self.temp_files = []
        self.template_path = self._create_graph_template()

        # Realistic context for a software architecture diagram
        self.context = {
            "project": "E-Commerce Platform",
            "version": "v2.0",
            "environment": "Production",
        }

    def tearDown(self):
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def _create_graph_template(self):
        """Create a graph template with layout."""
        prs = Presentation()

        # Create a graph layout slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add layout tag
        layout_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(4), Inches(0.5)
        )
        layout_box.text_frame.text = "% layout graph %"

        # Add title placeholder
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0), Inches(8), Inches(0.5)
        )
        title_box.text_frame.text = "{{ project }} - {{ version }}"

        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def test_software_architecture_diagram(self):
        """Test creating a realistic software architecture diagram."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "frontend",
                            "name": "Frontend",
                            "detail": "React.js Application",
                            "position": {"x": 1, "y": 2},
                        },
                        {
                            "id": "api_gateway",
                            "name": "API Gateway",
                            "detail": "Kong Gateway",
                            "position": {"x": 4, "y": 2},
                        },
                        {
                            "id": "auth_service",
                            "name": "Auth Service",
                            "detail": "OAuth2 + JWT",
                            "position": {"x": 7, "y": 1},
                        },
                        {
                            "id": "product_service",
                            "name": "Product Service",
                            "detail": "Microservice",
                            "position": {"x": 7, "y": 3},
                        },
                        {
                            "id": "database",
                            "name": "Database",
                            "detail": "PostgreSQL",
                            "position": {"x": 10, "y": 2},
                        },
                    ],
                    "edges": [
                        {"from": "frontend", "to": "api_gateway", "label": "HTTPS"},
                        {
                            "from": "api_gateway",
                            "to": "auth_service",
                            "label": "Authenticate",
                        },
                        {
                            "from": "api_gateway",
                            "to": "product_service",
                            "label": "Query",
                        },
                        {"from": "auth_service", "to": "database", "label": "User Data"},
                        {
                            "from": "product_service",
                            "to": "database",
                            "label": "Product Data",
                        },
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        self.assertTrue(os.path.exists(output_file))

        # Verify output structure
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

        # Verify slide has shapes
        slide = prs.slides[0]
        self.assertGreater(len(slide.shapes), 5)  # At least 5 nodes

        # Verify template variables were processed
        slide_text = self._extract_slide_text(slide)
        self.assertIn("E-Commerce Platform", slide_text)
        self.assertIn("v2.0", slide_text)

    def test_workflow_diagram(self):
        """Test creating a workflow/process diagram."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "start",
                            "name": "Start",
                            "detail": "Order Received",
                            "position": {"x": 1, "y": 2},
                        },
                        {
                            "id": "validate",
                            "name": "Validate",
                            "detail": "Check inventory",
                            "position": {"x": 3, "y": 2},
                        },
                        {
                            "id": "process",
                            "name": "Process Payment",
                            "detail": "Charge customer",
                            "position": {"x": 5, "y": 2},
                        },
                        {
                            "id": "fulfill",
                            "name": "Fulfill Order",
                            "detail": "Ship products",
                            "position": {"x": 7, "y": 2},
                        },
                        {
                            "id": "complete",
                            "name": "Complete",
                            "detail": "Order shipped",
                            "position": {"x": 9, "y": 2},
                        },
                    ],
                    "edges": [
                        {"from": "start", "to": "validate", "label": "Submit"},
                        {"from": "validate", "to": "process", "label": "Valid"},
                        {"from": "process", "to": "fulfill", "label": "Success"},
                        {"from": "fulfill", "to": "complete", "label": "Done"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

    def test_hierarchical_org_chart(self):
        """Test creating a hierarchical organization chart."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "ceo",
                            "name": "CEO",
                            "detail": "Chief Executive Officer",
                            "position": {"x": 4, "y": 1},
                        },
                        {
                            "id": "cto",
                            "name": "CTO",
                            "detail": "Chief Technology Officer",
                            "position": {"x": 2, "y": 3},
                        },
                        {
                            "id": "cfo",
                            "name": "CFO",
                            "detail": "Chief Financial Officer",
                            "position": {"x": 6, "y": 3},
                        },
                        {
                            "id": "eng_mgr",
                            "name": "Engineering Manager",
                            "detail": "Software Development",
                            "position": {"x": 1, "y": 5},
                        },
                        {
                            "id": "devops_mgr",
                            "name": "DevOps Manager",
                            "detail": "Infrastructure",
                            "position": {"x": 3, "y": 5},
                        },
                    ],
                    "edges": [
                        {"from": "ceo", "to": "cto", "label": "Reports to"},
                        {"from": "ceo", "to": "cfo", "label": "Reports to"},
                        {"from": "cto", "to": "eng_mgr", "label": "Manages"},
                        {"from": "cto", "to": "devops_mgr", "label": "Manages"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_multiple_architecture_diagrams(self):
        """Test creating multiple architecture diagrams in one presentation."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "dev_frontend",
                            "name": "Dev Frontend",
                            "detail": "localhost:3000",
                            "position": {"x": 1, "y": 2},
                        },
                        {
                            "id": "dev_backend",
                            "name": "Dev Backend",
                            "detail": "localhost:8000",
                            "position": {"x": 4, "y": 2},
                        },
                    ],
                    "edges": [
                        {"from": "dev_frontend", "to": "dev_backend", "label": "HTTP"},
                    ],
                },
            },
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "prod_lb",
                            "name": "Load Balancer",
                            "detail": "AWS ALB",
                            "position": {"x": 1, "y": 2},
                        },
                        {
                            "id": "prod_app",
                            "name": "Application Server",
                            "detail": "ECS Cluster",
                            "position": {"x": 4, "y": 2},
                        },
                        {
                            "id": "prod_db",
                            "name": "Database",
                            "detail": "RDS PostgreSQL",
                            "position": {"x": 7, "y": 2},
                        },
                    ],
                    "edges": [
                        {"from": "prod_lb", "to": "prod_app", "label": "HTTPS"},
                        {"from": "prod_app", "to": "prod_db", "label": "SQL"},
                    ],
                },
            },
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Verify we have 2 slides
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 2)

    def test_complex_network_topology(self):
        """Test creating a complex network topology with many nodes and edges."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "router",
                            "name": "Router",
                            "detail": "Main Gateway",
                            "position": {"x": 4, "y": 1},
                        },
                        {
                            "id": "switch1",
                            "name": "Switch 1",
                            "detail": "Floor 1",
                            "position": {"x": 2, "y": 3},
                        },
                        {
                            "id": "switch2",
                            "name": "Switch 2",
                            "detail": "Floor 2",
                            "position": {"x": 6, "y": 3},
                        },
                        {
                            "id": "server1",
                            "name": "Server 1",
                            "detail": "Web Server",
                            "position": {"x": 1, "y": 5},
                        },
                        {
                            "id": "server2",
                            "name": "Server 2",
                            "detail": "DB Server",
                            "position": {"x": 3, "y": 5},
                        },
                        {
                            "id": "server3",
                            "name": "Server 3",
                            "detail": "App Server",
                            "position": {"x": 5, "y": 5},
                        },
                        {
                            "id": "server4",
                            "name": "Server 4",
                            "detail": "Cache Server",
                            "position": {"x": 7, "y": 5},
                        },
                    ],
                    "edges": [
                        {"from": "router", "to": "switch1", "label": "1Gbps"},
                        {"from": "router", "to": "switch2", "label": "1Gbps"},
                        {"from": "switch1", "to": "server1", "label": "100Mbps"},
                        {"from": "switch1", "to": "server2", "label": "100Mbps"},
                        {"from": "switch2", "to": "server3", "label": "100Mbps"},
                        {"from": "switch2", "to": "server4", "label": "100Mbps"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

        # Verify we have many shapes
        slide = prs.slides[0]
        self.assertGreater(len(slide.shapes), 7)  # At least 7 nodes

    def test_graph_with_large_coordinates(self):
        """Test that slides auto-expand for large coordinate values.

        Note: Positions are in pixels, so we use large pixel values
        to test slide expansion.
        """
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "n1", "name": "Node 1", "position": {"x": 100, "y": 100}},
                        {"id": "n2", "name": "Node 2", "position": {"x": 1200, "y": 100}},
                        {"id": "n3", "name": "Node 3", "position": {"x": 100, "y": 800}},
                        {"id": "n4", "name": "Node 4", "position": {"x": 1200, "y": 800}},
                    ],
                    "edges": [
                        {"from": "n1", "to": "n2"},
                        {"from": "n1", "to": "n3"},
                        {"from": "n2", "to": "n4"},
                        {"from": "n3", "to": "n4"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Verify slide was expanded
        prs = Presentation(output_file)
        slide_width_inches = prs.slide_width / 914400  # Convert EMUs to inches
        slide_height_inches = prs.slide_height / 914400

        # Should be larger than default (positions are in pixels, converted to inches)
        # x=1200 pixels = 12.5 inches, plus node width (2.5) and margin (1) = 16 inches
        self.assertGreater(slide_width_inches, 10)
        # y=800 pixels = 8.33 inches, plus node height (1.5) and margin (1) = 10.83 inches
        self.assertGreater(slide_height_inches, 7.5)

    def test_data_flow_diagram(self):
        """Test creating a data flow diagram."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "user",
                            "name": "User",
                            "detail": "Web Browser",
                            "position": {"x": 1, "y": 3},
                        },
                        {
                            "id": "cdn",
                            "name": "CDN",
                            "detail": "CloudFront",
                            "position": {"x": 3, "y": 3},
                        },
                        {
                            "id": "app",
                            "name": "Application",
                            "detail": "Node.js",
                            "position": {"x": 5, "y": 3},
                        },
                        {
                            "id": "cache",
                            "name": "Cache",
                            "detail": "Redis",
                            "position": {"x": 7, "y": 2},
                        },
                        {
                            "id": "db",
                            "name": "Database",
                            "detail": "MongoDB",
                            "position": {"x": 7, "y": 4},
                        },
                    ],
                    "edges": [
                        {"from": "user", "to": "cdn", "label": "Request"},
                        {"from": "cdn", "to": "app", "label": "Forward"},
                        {"from": "app", "to": "cache", "label": "Check"},
                        {"from": "app", "to": "db", "label": "Query"},
                        {"from": "cache", "to": "app", "label": "Data"},
                        {"from": "db", "to": "app", "label": "Data"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_graph_with_nested_nodes_placeholder(self):
        """Test graph with nested nodes (placeholder functionality)."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "container",
                            "name": "Container",
                            "detail": "Docker Container",
                            "position": {"x": 1, "y": 1},
                        },
                        {
                            "id": "service1",
                            "name": "Service 1",
                            "detail": "Web Service",
                            "position": {"x": 2, "y": 2},
                            "parent": "container",
                        },
                        {
                            "id": "service2",
                            "name": "Service 2",
                            "detail": "API Service",
                            "position": {"x": 2, "y": 3},
                            "parent": "container",
                        },
                    ],
                    "edges": [
                        {"from": "service1", "to": "service2", "label": "Internal"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        # Should succeed even though parent functionality is placeholder
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def _extract_slide_text(self, slide):
        """Extract all text content from a slide."""
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
            elif hasattr(shape, "text_frame"):
                text_content.append(shape.text_frame.text)
        return " ".join(text_content)


if __name__ == "__main__":
    unittest.main()

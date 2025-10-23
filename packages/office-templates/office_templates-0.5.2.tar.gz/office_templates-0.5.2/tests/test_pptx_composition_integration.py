import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import compose_pptx


class TestPPTXCompositionIntegration(unittest.TestCase):
    """Integration tests for PPTX composition with realistic scenarios."""

    def setUp(self):
        self.temp_files = []
        self.template1_path = self._create_presentation_template()
        self.template2_path = self._create_data_template()

        # Realistic context
        self.context = {
            "company": "Acme Corp",
            "quarter": "Q1 2024",
            "presenter": {"name": "John Smith", "title": "VP of Sales"},
            "sales_data": [
                {"month": "January", "revenue": 150000, "target": 140000},
                {"month": "February", "revenue": 165000, "target": 150000},
                {"month": "March", "revenue": 180000, "target": 160000},
            ],
            "data_count": 3,  # Length of sales_data for template usage
            "top_products": ["Product A", "Product B", "Product C"],
            "product_list": "Product A, Product B, Product C",  # Joined version for template
        }

    def tearDown(self):
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def _create_presentation_template(self):
        """Create a presentation template with various slide layouts."""
        prs = Presentation()

        # Title slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        slide1.shapes.title.text = "Company Quarterly Report"
        if len(slide1.shapes) > 1:
            slide1.shapes[1].text_frame.text = "Prepared by Sales Team"

        # Content slide with tagged layout
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Executive Summary"
        
        # Add %layout% tag in its own shape
        layout_box = slide2.shapes.add_textbox(
            Inches(1), Inches(2), Inches(6), Inches(0.5)
        )
        layout_box.text_frame.text = "% layout summary %"
        
        # Add template variable in separate shape
        content_box = slide2.shapes.add_textbox(
            Inches(1), Inches(3), Inches(6), Inches(1)
        )
        content_box.text_frame.text = "{{ summary_content }}"

        # Chart slide
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        title_box = slide3.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_box.text_frame.text = "Sales Performance Chart"
        
        # Add %layout% tag in its own shape
        layout_box = slide3.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(0.5)
        )
        layout_box.text_frame.text = "% layout chart %"
        
        # Add template variable in separate shape
        chart_placeholder = slide3.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(4)
        )
        chart_placeholder.text_frame.text = "{{ chart_title }}"

        # Save template
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def _create_data_template(self):
        """Create a data-focused template with tables and lists."""
        prs = Presentation()

        # Data table slide
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
        title_box = slide1.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )
        title_box.text_frame.text = "Monthly Data"
        
        # Add %layout% tag in its own shape
        layout_box = slide1.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8), Inches(0.5)
        )
        layout_box.text_frame.text = "% layout data_table %"
        
        # Add template variable in separate shape
        table_marker = slide1.shapes.add_textbox(
            Inches(1), Inches(2), Inches(8), Inches(3)
        )
        table_marker.text_frame.text = "{{ table_title }}"

        # Product list slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Top Products"
        if len(slide2.shapes) > 1:
            slide2.shapes[1].text_frame.text = "Our best-selling products this quarter\n{{ product_summary }}"

        # Save template
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def test_full_presentation_composition(self):
        """Test creating a complete presentation using multiple templates and layout types."""
        slide_specs = [
            # Title slide using placeholders
            {
                "layout": "Company Quarterly Report",
                "placeholders": [
                    "{{ company }} {{ quarter }} Report",
                    "Presented by {{ presenter.name }}, {{ presenter.title }}",
                ],
            },
            # Executive summary using tagged layout
            {
                "layout": "summary",
                "summary_content": "This quarter exceeded expectations with {{ data_count }} months of strong performance.",
            },
            # Data table using tagged layout from second template
            {
                "layout": "data_table",
                "table_title": "{{ quarter }} Monthly Performance",
            },
            # Chart slide using tagged layout
            {
                "layout": "chart",
                "chart_title": "Revenue vs Target - {{ quarter }}",
            },
            # Product list using slide title as layout
            {
                "layout": "Top Products",
                "product_summary": "Featured products: {{ product_list }}",
            },
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path, self.template2_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
            use_all_slides_as_layouts_by_title=True,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Verify output
        self.assertTrue(os.path.exists(output_file))
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 5)

        # Check that slides have content and verify specific processed content
        for slide in prs.slides:
            self.assertGreater(
                len(slide.shapes), 0, "Each slide should have at least one shape"
            )

        # Verify specific slide content was processed correctly
        # Slide 1: Title slide with placeholders
        slide1_text = self._extract_slide_text(prs.slides[0])
        # Verify specific slide content was processed correctly
        # Slide 1: Title slide with placeholders
        slide1_text = self._extract_slide_text(prs.slides[0])
        self.assertIn("Acme Corp Q1 2024 Report", slide1_text)
        self.assertIn("John Smith, VP of Sales", slide1_text)

        # Slide 2: Executive summary with template processing
        slide2_text = self._extract_slide_text(prs.slides[1])
        self.assertIn("3 months of strong performance", slide2_text)  # {{ data_count }} processed

        # Slide 3: Data table with processed title
        slide3_text = self._extract_slide_text(prs.slides[2])
        self.assertIn("Q1 2024 Monthly Performance", slide3_text)

        # Slide 4: Chart with processed title
        slide4_text = self._extract_slide_text(prs.slides[3])
        self.assertIn("Revenue vs Target - Q1 2024", slide4_text)

        # Slide 5: Product list with joined products
        slide5_text = self._extract_slide_text(prs.slides[4])
        self.assertIn("Product A, Product B, Product C", slide5_text)

        # Verify that all %layout% shapes have been removed from the output
        for i, slide in enumerate(prs.slides):
            layout_shapes = []
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'text'):
                    text = shape.text_frame.text.strip()
                    if "% layout" in text and "%" in text:
                        layout_shapes.append(shape)
            self.assertEqual(len(layout_shapes), 0, 
                           f"Found %layout% shapes in output slide {i+1}: {[s.text_frame.text for s in layout_shapes]}")

    def _extract_slide_text(self, slide):
        """Extract all text content from a slide for validation."""
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text_content.append(shape.text)
            elif hasattr(shape, 'text_frame'):
                text_content.append(shape.text_frame.text)
        return " ".join(text_content)

    def test_mixed_layout_types_composition(self):
        """Test composition using all three layout discovery methods."""
        slide_specs = [
            # Using master layout (slide layouts)
            {"layout": "Title Slide", "content": "Master layout test"},
            # Using tagged layout
            {"layout": "summary", "summary_content": "Tagged layout test"},
            # Using title-based layout
            {"layout": "Top Products", "product_summary": "Title-based layout test"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path, self.template2_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
            use_all_slides_as_layouts_by_title=True,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Verify output
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 3)

    def test_template_processing_with_context(self):
        """Test that template tags are properly processed in composed slides."""
        slide_specs = [
            {
                "layout": "Company Quarterly Report",
                "placeholders": ["{{ company }} Report", "By {{ presenter.name }}"],
                "custom_field": "Revenue grew {{ data_count }}x this quarter",
            },
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_all_slides_as_layouts_by_title=True,
        )

        # Should succeed
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check that placeholders were processed
        prs = Presentation(output_file)
        slide = prs.slides[0]

        # Find placeholder shapes and verify they contain processed content
        placeholder_found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                text = shape.text_frame.text
                if "Acme Corp Report" in text or "By John Smith" in text:
                    placeholder_found = True
                    break

        self.assertTrue(placeholder_found, "Placeholder processing should have occurred")

    def test_error_handling_with_complex_scenarios(self):
        """Test error handling in complex composition scenarios."""
        slide_specs = [
            # Valid slide
            {"layout": "summary", "summary_content": "Valid content"},
            # Invalid layout
            {"layout": "NonExistentLayout", "content": "This should fail"},
            # Valid slide after error
            {"layout": "Top Products", "product_summary": "This should also be processed"},
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template1_path, self.template2_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
            use_all_slides_as_layouts_by_title=True,
        )

        # Should fail due to invalid layout
        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(
            any("Layout 'NonExistentLayout' not found" in error for error in errors)
        )


if __name__ == "__main__":
    unittest.main()

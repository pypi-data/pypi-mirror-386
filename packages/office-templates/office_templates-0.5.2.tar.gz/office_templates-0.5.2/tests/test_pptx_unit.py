import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import render_pptx

from tests.utils import has_view_permission


# Dummy objects for testing renderer behavior.
class DummyUser:
    def __init__(self, name, email):
        self.name = name
        self.email = email
        self._meta = True  # Simulate a Django model

    def __str__(self):
        return self.name


class DummyRequestUser:
    def has_perm(self, perm, obj):
        # Deny permission if object's name contains "deny"
        if hasattr(obj, "name") and "deny" in obj.name.lower():
            return False
        return True


class TestRendererUnit(unittest.TestCase):
    def setUp(self):
        # Create a minimal PPTX file with one slide and one text box.
        self.prs = Presentation()
        slide_layout = self.prs.slide_layouts[5]  # use a blank layout
        self.slide = self.prs.slides.add_slide(slide_layout)
        self.textbox = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(1)
        )
        self.textframe = self.textbox.text_frame
        self.shape_index = len(self.slide.shapes) - 1
        # Insert a paragraph with a pure placeholder.
        p = self.textframe.paragraphs[0]
        p.text = "Hello, {{ user.name }}"

        # Create a temporary file to save the PPTX.
        self.temp_input = tempfile.mktemp(suffix=".pptx")
        self.temp_output = tempfile.mktemp(suffix=".pptx")
        self.prs.save(self.temp_input)

        # Prepare dummy context.
        self.user = DummyUser("Alice", "alice@example.com")
        self.context = {
            "user": self.user,
        }
        self.request_user = DummyRequestUser()

    def tearDown(self):
        # Clean up temporary files.
        if os.path.exists(self.temp_input):
            os.remove(self.temp_input)
        if os.path.exists(self.temp_output):
            os.remove(self.temp_output)

    def test_textbox_pure_placeholder_normal_mode(self):
        # Test that pure placeholder in a text box is rendered correctly.
        rendered, _ = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            None,
        )
        # Load output and check text from slide[0], shape[0].
        prs_out = Presentation(rendered)
        shape = prs_out.slides[0].shapes[self.shape_index]
        txt = shape.text_frame.text
        self.assertEqual(txt, "Hello, Alice")

    def test_mixed_text_list_joining(self):
        # Test mixed text where placeholder resolves to a list.
        # For example: "Emails: {{ user.email }}"
        # We'll set user.email to be a list for this test.
        self.user.email = ["alice@example.com", "alice.alt@example.com"]
        self.textframe.paragraphs[0].text = "Emails: {{ user.email }}"
        self.prs.save(self.temp_input)
        rendered, _ = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            None,
        )
        prs_out = Presentation(rendered)
        shape = prs_out.slides[0].shapes[self.shape_index]
        text = shape.text_frame.text
        self.assertEqual(text, "Emails: alice@example.com, alice.alt@example.com")

    def test_permission_denied_pure(self):
        self.user.name = "DenyUser"
        # Retrieve entire user (Django-like object), forcing permission denial:
        self.textframe.paragraphs[0].text = "{{ user }}"
        self.prs.save(self.temp_input)
        check_permissions = lambda obj: has_view_permission(obj, self.request_user)
        _, errors = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            check_permissions,
        )
        self.assertTrue(len(errors) == 1)


if __name__ == "__main__":
    unittest.main()

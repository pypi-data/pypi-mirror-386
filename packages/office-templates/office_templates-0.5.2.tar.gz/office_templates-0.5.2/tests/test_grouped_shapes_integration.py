"""
Integration tests for grouped shapes functionality in PPTX rendering.
"""

import os
import tempfile
import unittest

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from office_templates.office_renderer.pptx.render import render_pptx


class TestGroupedShapesIntegration(unittest.TestCase):
    def setUp(self):
        """Create test presentations with grouped shapes."""
        self.temp_files = []
        self.template_path = self._create_test_template_with_groups()
        
        # Test context
        self.context = {
            "title": "Test Presentation",
            "grouped_text": "Rendered grouped variable",
            "another_var": "Rendered another variable",
            "chart_data": [
                ["Month", "Sales"],
                ["Jan", 100],
                ["Feb", 150],
                ["Mar", 200],
            ],
        }

    def tearDown(self):
        """Clean up temporary files."""
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def _create_test_template_with_groups(self):
        """Create a test template with grouped shapes."""
        prs = Presentation()
        
        # Slide 1: Text shapes in a group
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        
        text_box1 = slide1.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        text_box1.text_frame.text = "Title: {{ title }}"
        
        text_box2 = slide1.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(3), Inches(1)
        )
        text_box2.text_frame.text = "Content: {{ grouped_text }}"
        
        text_box3 = slide1.shapes.add_textbox(
            Inches(1), Inches(4), Inches(3), Inches(1)
        )
        text_box3.text_frame.text = "Footer: {{ another_var }}"
        
        # Group the text boxes
        group1 = slide1.shapes.add_group_shape([text_box1, text_box2, text_box3])
        
        # Slide 2: Mixed content in a group (text + table would be ideal, but let's start simple)
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Add some ungrouped content for comparison
        ungrouped_text = slide2.shapes.add_textbox(
            Inches(5), Inches(1), Inches(2), Inches(1)
        )
        ungrouped_text.text_frame.text = "Ungrouped: {{ title }}"
        
        # Add grouped content  
        grouped_text1 = slide2.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        grouped_text1.text_frame.text = "In group 1: {{ grouped_text }}"
        
        grouped_text2 = slide2.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(3), Inches(1)
        )
        grouped_text2.text_frame.text = "In group 2: {{ another_var }}"
        
        group2 = slide2.shapes.add_group_shape([grouped_text1, grouped_text2])
        
        # Save template
        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def test_render_grouped_text_shapes(self):
        """Test that text shapes within groups are properly rendered."""
        # Create output file
        output_path = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_path)
        
        # Render the template
        result, errors = render_pptx(
            template=self.template_path,
            context=self.context,
            output=output_path,
            check_permissions=None,
        )
        
        # Should complete without errors now that we support grouped shapes
        self.assertIsNone(errors, f"Rendering should succeed: {errors}")
        
        # Load the rendered presentation
        if os.path.exists(output_path):
            rendered_prs = Presentation(output_path)
            
            # Check slide 1 - should have rendered grouped content
            slide1 = rendered_prs.slides[0]
            self._verify_grouped_shapes_rendered(slide1)
            
            # Check slide 2 - should have both grouped and ungrouped content rendered
            slide2 = rendered_prs.slides[1]
            self._verify_mixed_content_rendered(slide2)

    def _verify_grouped_shapes_rendered(self, slide):
        """Verify that grouped shapes were properly rendered."""
        group_found = False
        grouped_texts = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_found = True
                for grouped_shape in shape.shapes:
                    if hasattr(grouped_shape, "text_frame") and grouped_shape.text_frame:
                        grouped_texts.append(grouped_shape.text_frame.text)
        
        self.assertTrue(group_found, "Should find grouped shapes on slide")
        
        # These should now work with grouped shapes support
        expected_texts = [
            "Title: Test Presentation",
            "Content: Rendered grouped variable", 
            "Footer: Rendered another variable"
        ]
        
        for expected_text in expected_texts:
            self.assertIn(expected_text, grouped_texts, 
                         f"Expected text '{expected_text}' not found in grouped texts: {grouped_texts}")

    def _verify_mixed_content_rendered(self, slide):
        """Verify that both grouped and ungrouped content was rendered."""
        ungrouped_texts = []
        grouped_texts = []
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for grouped_shape in shape.shapes:
                    if hasattr(grouped_shape, "text_frame") and grouped_shape.text_frame:
                        grouped_texts.append(grouped_shape.text_frame.text)
            else:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    ungrouped_texts.append(shape.text_frame.text)
        
        # Ungrouped content should be rendered correctly (this works currently)
        self.assertIn("Ungrouped: Test Presentation", ungrouped_texts)
        
        # Grouped content should be rendered correctly now
        expected_grouped = [
            "In group 1: Rendered grouped variable",
            "In group 2: Rendered another variable"
        ]
        
        for expected_text in expected_grouped:
            self.assertIn(expected_text, grouped_texts,
                         f"Expected grouped text '{expected_text}' not found: {grouped_texts}")

    def test_render_template_with_no_groups(self):
        """Test that templates without groups still work correctly (regression test)."""
        # Create a simple template without groups
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        text_box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        text_box.text_frame.text = "Simple: {{ title }}"
        
        template_path = tempfile.mktemp(suffix=".pptx")
        prs.save(template_path)
        self.temp_files.append(template_path)
        
        # Render
        output_path = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_path)
        
        result, errors = render_pptx(
            template=template_path,
            context=self.context,
            output=output_path,
            check_permissions=None,
        )
        
        # Should work fine
        self.assertIsNone(errors, f"Should render without errors: {errors}")
        
        # Verify content was rendered
        rendered_prs = Presentation(output_path)
        slide = rendered_prs.slides[0]
        
        text_found = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                if "Simple: Test Presentation" in shape.text_frame.text:
                    text_found = True
                    break
        
        self.assertTrue(text_found, "Should find rendered text in non-grouped template")


if __name__ == "__main__":
    unittest.main()
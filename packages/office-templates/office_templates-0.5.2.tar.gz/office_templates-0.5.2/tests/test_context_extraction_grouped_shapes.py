"""
Test context extraction with grouped shapes.
"""

import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer.context_extractor import extract_context_keys_from_pptx


class TestContextExtractionGroupedShapes(unittest.TestCase):
    def setUp(self):
        """Create a test presentation with grouped shapes containing context variables."""
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        
        # Create ungrouped shape for comparison
        ungrouped_text = self.slide.shapes.add_textbox(
            Inches(5), Inches(1), Inches(2), Inches(1)
        )
        ungrouped_text.text_frame.text = "Ungrouped: {{ ungrouped_var }}"
        
        # Create grouped shapes with various variable types
        text_box1 = self.slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        text_box1.text_frame.text = "Simple grouped: {{ simple_grouped }}"
        
        text_box2 = self.slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(3), Inches(1)
        )
        text_box2.text_frame.text = "Object grouped: {{ user.name }} and {{ user.email }}"
        
        text_box3 = self.slide.shapes.add_textbox(
            Inches(1), Inches(4), Inches(3), Inches(1)
        )
        text_box3.text_frame.text = "Array access: {{ items[0].title }}"
        
        # Group the shapes
        self.group = self.slide.shapes.add_group_shape([text_box1, text_box2, text_box3])

    def test_extract_context_keys_from_grouped_shapes(self):
        """Test that context extraction works with grouped shapes."""
        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            self.prs.save(temp_file.name)
            
            # Extract context keys
            result = extract_context_keys_from_pptx(temp_file.name)
            
            # Check simple fields
            simple_fields = result["simple_fields"]
            self.assertIn("simple_grouped", simple_fields)
            self.assertIn("ungrouped_var", simple_fields)
            
            # Check object fields  
            object_fields = result["object_fields"]
            self.assertIn("user", object_fields)
            self.assertIn("items", object_fields)
            
            print(f"Extracted simple fields: {simple_fields}")
            print(f"Extracted object fields: {object_fields}")

    def test_extract_context_keys_handles_both_grouped_and_ungrouped(self):
        """Test that both grouped and ungrouped shapes are processed."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            self.prs.save(temp_file.name)
            
            result = extract_context_keys_from_pptx(temp_file.name)
            
            # Should find variables from both grouped and ungrouped shapes
            simple_fields = result["simple_fields"]
            
            # From ungrouped shape
            self.assertIn("ungrouped_var", simple_fields)
            
            # From grouped shapes
            self.assertIn("simple_grouped", simple_fields)

    def test_extract_context_keys_with_nested_variables(self):
        """Test extraction of complex nested variables in grouped shapes."""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            self.prs.save(temp_file.name)
            
            result = extract_context_keys_from_pptx(temp_file.name)
            
            object_fields = result["object_fields"]
            
            # Should extract top-level keys from complex expressions
            self.assertIn("user", object_fields)  # From {{ user.name }} and {{ user.email }}
            self.assertIn("items", object_fields)  # From {{ items[0].title }}


if __name__ == "__main__":
    unittest.main()
import tempfile
import unittest
import os
from openpyxl import Workbook
from pptx import Presentation

from office_templates.office_renderer.context_extractor import (
    extract_context_keys,
    extract_context_keys_from_xlsx,
    extract_context_keys_from_pptx,
    extract_top_level_context_keys_from_text,
)
from office_templates.office_renderer.utils import identify_file_type


class TestContextExtractor(unittest.TestCase):
    def test_extract_top_level_context_keys_from_text(self):
        """Test extraction of context keys from text strings."""
        # Simple field
        result = extract_top_level_context_keys_from_text("{{ simple_field }}")
        self.assertEqual(result["simple_fields"], ["simple_field"])
        self.assertEqual(result["object_fields"], [])

        # Object field with dot notation
        result = extract_top_level_context_keys_from_text("{{ user.name }}")
        self.assertEqual(result["simple_fields"], [])
        self.assertEqual(result["object_fields"], ["user"])

        # Object field with bracket notation
        result = extract_top_level_context_keys_from_text("{{ users[0] }}")
        self.assertEqual(result["simple_fields"], [])
        self.assertEqual(result["object_fields"], ["users"])

        # Mixed content
        text = "Hello {{ user.name }}, you work at {{ company }} in {{ department.name }}"
        result = extract_top_level_context_keys_from_text(text)
        self.assertEqual(result["simple_fields"], ["company"])
        self.assertEqual(result["object_fields"], ["department", "user"])

        # Ignored keys
        result = extract_top_level_context_keys_from_text("{{ now }} {{ loop_count }}")
        self.assertEqual(result["simple_fields"], [])
        self.assertEqual(result["object_fields"], [])

    def test_extract_context_keys_from_xlsx(self):
        """Test context key extraction from XLSX files."""
        # Create a temporary XLSX file
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            wb = Workbook()
            ws = wb.active
            ws.title = "Test Sheet"
            
            # Add various types of placeholders
            ws['A1'] = '{{ user.name }}'
            ws['B1'] = '{{ company }}'
            ws['A2'] = '{{ user.email }}'
            ws['B2'] = '{{ department.budget }}'
            ws['A3'] = '{{ simple_field }}'
            ws['B3'] = '{{ user.address.city }}'
            
            # Add a second worksheet
            ws2 = wb.create_sheet("Second Sheet")
            ws2['A1'] = '{{ another_field }}'
            ws2['B1'] = '{{ nested.object.value }}'
            
            wb.save(tmp.name)
            tmp_path = tmp.name

        try:
            # Test with file path
            result = extract_context_keys_from_xlsx(tmp_path)
            
            expected_simple = ['another_field', 'company', 'simple_field']
            expected_object = ['department', 'nested', 'user']
            
            self.assertEqual(result['simple_fields'], expected_simple)
            self.assertEqual(result['object_fields'], expected_object)

            # Test with file-like object
            with open(tmp_path, 'rb') as f:
                result = extract_context_keys_from_xlsx(f)
                self.assertEqual(result['simple_fields'], expected_simple)
                self.assertEqual(result['object_fields'], expected_object)

        finally:
            os.unlink(tmp_path)

    def test_extract_context_keys_from_pptx(self):
        """Test context key extraction from PPTX files."""
        # Create a temporary PPTX file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            prs = Presentation()
            
            # Add a slide with content
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
            title = slide.shapes.title
            title.text = "Test Slide {{ user.name }}"
            
            content = slide.placeholders[1]
            content.text = "Hello {{ company }}, your budget is {{ department.budget }}"
            
            # Add another slide
            slide2 = prs.slides.add_slide(prs.slide_layouts[1])
            slide2.shapes.title.text = "Second slide {{ simple_field }}"
            slide2.placeholders[1].text = "Address: {{ user.address.city }}"
            
            prs.save(tmp.name)
            tmp_path = tmp.name

        try:
            # Test with file path
            result = extract_context_keys_from_pptx(tmp_path)
            
            self.assertIn('company', result['simple_fields'])
            self.assertIn('simple_field', result['simple_fields'])
            self.assertIn('user', result['object_fields'])
            self.assertIn('department', result['object_fields'])

            # Test with file-like object
            with open(tmp_path, 'rb') as f:
                result = extract_context_keys_from_pptx(f)
                self.assertIn('company', result['simple_fields'])
                self.assertIn('user', result['object_fields'])

        finally:
            os.unlink(tmp_path)

    def test_extract_context_keys_file_type_detection(self):
        """Test that extract_context_keys properly detects file types."""
        # Test with XLSX
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp_xlsx:
            wb = Workbook()
            ws = wb.active
            ws['A1'] = '{{ test_field }}'
            wb.save(tmp_xlsx.name)
            xlsx_path = tmp_xlsx.name

        # Test with PPTX
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "{{ test_field }}"
            prs.save(tmp_pptx.name)
            pptx_path = tmp_pptx.name

        try:
            # Test XLSX file type detection
            xlsx_result = extract_context_keys(xlsx_path)
            self.assertIn('test_field', xlsx_result['simple_fields'])

            # Test PPTX file type detection
            pptx_result = extract_context_keys(pptx_path)
            self.assertIn('test_field', pptx_result['simple_fields'])

            # Verify file type identification works
            self.assertEqual(identify_file_type(xlsx_path), "xlsx")
            self.assertEqual(identify_file_type(pptx_path), "pptx")

        finally:
            os.unlink(xlsx_path)
            os.unlink(pptx_path)

    def test_extract_context_keys_empty_files(self):
        """Test behavior with empty files."""
        # Test with empty XLSX
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp_xlsx:
            wb = Workbook()
            wb.save(tmp_xlsx.name)
            xlsx_path = tmp_xlsx.name

        # Test with empty PPTX
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
            prs = Presentation()
            prs.save(tmp_pptx.name)
            pptx_path = tmp_pptx.name

        try:
            # Empty files should return empty results
            xlsx_result = extract_context_keys(xlsx_path)
            self.assertEqual(xlsx_result['simple_fields'], [])
            self.assertEqual(xlsx_result['object_fields'], [])

            pptx_result = extract_context_keys(pptx_path)
            self.assertEqual(pptx_result['simple_fields'], [])
            self.assertEqual(pptx_result['object_fields'], [])

        finally:
            os.unlink(xlsx_path)
            os.unlink(pptx_path)


if __name__ == '__main__':
    unittest.main()
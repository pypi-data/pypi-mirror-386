import os
import tempfile
import unittest

from pptx import Presentation
from pptx.util import Inches

from office_templates.office_renderer import compose_pptx


class TestGraphsUnit(unittest.TestCase):
    """Unit tests for graph composition functionality."""

    def setUp(self):
        self.temp_files = []
        self.template_path = self._create_test_template()
        self.context = {
            "company": "Test Corp",
            "node_prefix": "Node",
        }

    def tearDown(self):
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def _create_test_template(self):
        """Create a simple test template."""
        prs = Presentation()

        # Add a blank slide as a graph layout
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add layout tag
        layout_box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(0.5)
        )
        layout_box.text_frame.text = "% layout graph %"

        temp_file = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_file)
        self.temp_files.append(temp_file)
        return temp_file

    def test_simple_graph_with_nodes_only(self):
        """Test creating a simple graph with just nodes, no edges."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "node1",
                            "name": "First Node",
                            "detail": "Details about first node",
                            "position": {"x": 1, "y": 1},
                        },
                        {
                            "id": "node2",
                            "name": "Second Node",
                            "detail": "Details about second node",
                            "position": {"x": 4, "y": 1},
                        },
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        self.assertTrue(os.path.exists(output_file))

        # Verify the output
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

        # Check that slide has shapes (nodes)
        slide = prs.slides[0]
        self.assertGreater(len(slide.shapes), 0)

    def test_graph_with_nodes_and_edges(self):
        """Test creating a graph with nodes and edges."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "A",
                            "name": "Node A",
                            "position": {"x": 1, "y": 2},
                        },
                        {
                            "id": "B",
                            "name": "Node B",
                            "position": {"x": 4, "y": 2},
                        },
                        {
                            "id": "C",
                            "name": "Node C",
                            "position": {"x": 7, "y": 2},
                        },
                    ],
                    "edges": [
                        {"from": "A", "to": "B", "label": "connects to"},
                        {"from": "B", "to": "C", "label": "flows to"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        self.assertTrue(os.path.exists(output_file))

        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

    def test_multiple_graphs_create_multiple_slides(self):
        """Test that multiple graphs create multiple slides."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "1", "name": "Node 1", "position": {"x": 1, "y": 1}},
                    ],
                    "edges": [],
                },
            },
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "2", "name": "Node 2", "position": {"x": 1, "y": 1}},
                    ],
                    "edges": [],
                },
            },
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "3", "name": "Node 3", "position": {"x": 1, "y": 1}},
                    ],
                    "edges": [],
                },
            },
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 3)

    def test_template_variable_processing_in_nodes(self):
        """Test that template variables in node names are processed."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "test",
                            "name": "{{ company }}",
                            "detail": "Prefix: {{ node_prefix }}",
                            "position": {"x": 1, "y": 1},
                        },
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

        # Check that variables were processed
        prs = Presentation(output_file)
        slide = prs.slides[0]

        # Find text containing processed variable
        found_company = False
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text
                if "Test Corp" in text:
                    found_company = True
                    break

        self.assertTrue(found_company, "Template variable should be processed")

    def test_edge_label_processing(self):
        """Test that edge labels with template variables are processed."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "A", "name": "Node A", "position": {"x": 1, "y": 1}},
                        {"id": "B", "name": "Node B", "position": {"x": 4, "y": 1}},
                    ],
                    "edges": [
                        {"from": "A", "to": "B", "label": "{{ company }} edge"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_node_without_position_error(self):
        """Test error handling when node is missing position."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "bad_node",
                            "name": "Bad Node",
                            # Missing position
                        },
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'position'" in error for error in errors))

    def test_node_without_id_error(self):
        """Test error handling when node is missing id."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            # Missing id
                            "name": "Node",
                            "position": {"x": 1, "y": 1},
                        },
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'id'" in error for error in errors))

    def test_node_without_name_error(self):
        """Test error handling when node is missing name."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "test",
                            # Missing name
                            "position": {"x": 1, "y": 1},
                        },
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'name'" in error for error in errors))

    def test_edge_with_unknown_node_error(self):
        """Test error handling when edge references unknown node."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "A", "name": "Node A", "position": {"x": 1, "y": 1}},
                    ],
                    "edges": [
                        {"from": "A", "to": "B"},  # B doesn't exist
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("unknown target node" in error for error in errors))

    def test_graph_missing_nodes_key_error(self):
        """Test error handling when graph is missing nodes key."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    # Missing nodes
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'nodes' key" in error for error in errors))

    def test_graph_missing_edges_key_error(self):
        """Test error handling when graph is missing edges key."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "test", "name": "Node", "position": {"x": 1, "y": 1}},
                    ],
                    # Missing edges
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("missing 'edges' key" in error for error in errors))

    def test_no_graphs_error(self):
        """Test error handling when no graphs are provided."""
        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=[],
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        self.assertTrue(any("No slides specified" in error for error in errors))

    def test_no_template_files_invalid_layout_error(self):
        """Test error handling when no template files and invalid layout used."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "test", "name": "Node", "position": {"x": 1, "y": 1}},
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNone(result)
        self.assertIsNotNone(errors)
        # Should fail because "graph" is not a default layout name
        self.assertTrue(any("Layout 'graph' not found" in error for error in errors))

    def test_node_with_parent_placeholder(self):
        """Test that nodes with parent key are accepted (placeholder functionality)."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "parent",
                            "name": "Parent Node",
                            "position": {"x": 1, "y": 1},
                        },
                        {
                            "id": "child",
                            "name": "Child Node",
                            "position": {"x": 2, "y": 2},
                            "parent": "parent",  # This is placeholder for now
                        },
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        # Should succeed - parent is accepted but not yet fully implemented
        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_node_without_detail(self):
        """Test that nodes without detail field work correctly."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "simple",
                            "name": "Simple Node",
                            # No detail field
                            "position": {"x": 1, "y": 1},
                        },
                    ],
                    "edges": [],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_edge_without_label(self):
        """Test that edges without label field work correctly."""
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {"id": "A", "name": "Node A", "position": {"x": 1, "y": 1}},
                        {"id": "B", "name": "Node B", "position": {"x": 4, "y": 1}},
                    ],
                    "edges": [
                        {"from": "A", "to": "B"},  # No label
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        self.assertIsNotNone(result)
        self.assertIsNone(errors)

    def test_pixel_based_positions_from_frontend(self):
        """Test that pixel-based positions from frontend graphs work correctly.

        This test uses the exact positions from the bug report where positions
        like x=390.99999928474426 were being interpreted as inches instead of pixels,
        causing an error: "value must be in range(914400, 51206400) (1-56 inches),
        got 360730799"
        """
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "d65d0243-a79a-4605-9df7-544522c20aa8",
                            "name": "Node 1",
                            "position": {"x": 20, "y": 33.91666666666667},
                        },
                        {
                            "id": "645f7aa1-5507-432a-8ad8-c61b860f3052",
                            "name": "Node 2",
                            "position": {"x": 189.5, "y": 33.91666666666667},
                        },
                        {
                            "id": "a6695bf7-d8f6-4baa-89b9-bf11b65379ae",
                            "name": "Node 3",
                            "position": {"x": 390.99999928474426, "y": 227.75},
                        },
                        {
                            "id": "111ede80-74cd-4c94-80e0-610d471eadcb",
                            "name": "Node 4",
                            "position": {"x": 390.99999928474426, "y": 287},
                        },
                        {
                            "id": "24466882-e684-4e41-8a5c-25fab255fece",
                            "name": "Node 5",
                            "position": {"x": 390.99999928474426, "y": 168.5},
                        },
                        {
                            "id": "ac781a02-f6b5-4986-95cc-27dbd5b1ff97",
                            "name": "Node 6",
                            "position": {"x": 370.99999928474426, "y": 20},
                        },
                    ],
                    "edges": [
                        {
                            "from": "d65d0243-a79a-4605-9df7-544522c20aa8",
                            "to": "a6695bf7-d8f6-4baa-89b9-bf11b65379ae",
                        },
                        {
                            "from": "645f7aa1-5507-432a-8ad8-c61b860f3052",
                            "to": "111ede80-74cd-4c94-80e0-610d471eadcb",
                        },
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        # Should succeed without errors
        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        self.assertTrue(os.path.exists(output_file))

        # Verify the output
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

        # Check that all nodes were created
        slide = prs.slides[0]
        # Should have 6 nodes + 2 connectors (and maybe layout box)
        self.assertGreater(len(slide.shapes), 6)

    def test_enormous_graph_auto_scaling(self):
        """Test that enormous graphs automatically scale to fit within PowerPoint limits.

        PowerPoint has a maximum slide dimension of 56 inches. For very large graphs
        that would exceed this limit, the system should automatically scale down
        positions, node sizes, and font sizes to fit.
        """
        slide_specs = [
            {
                "layout": "graph",
                "graph": {
                    "nodes": [
                        {
                            "id": "n1",
                            "name": "Node 1",
                            "detail": "Top Left",
                            "position": {"x": 100, "y": 100},
                        },
                        {
                            "id": "n2",
                            "name": "Node 2",
                            "detail": "Top Right",
                            "position": {"x": 10000, "y": 100},
                        },
                        {
                            "id": "n3",
                            "name": "Node 3",
                            "detail": "Bottom Left",
                            "position": {"x": 100, "y": 8000},
                        },
                        {
                            "id": "n4",
                            "name": "Node 4",
                            "detail": "Bottom Right",
                            "position": {"x": 10000, "y": 8000},
                        },
                    ],
                    "edges": [
                        {"from": "n1", "to": "n2", "label": "Top edge"},
                        {"from": "n1", "to": "n3", "label": "Left edge"},
                        {"from": "n2", "to": "n4", "label": "Right edge"},
                        {"from": "n3", "to": "n4", "label": "Bottom edge"},
                    ],
                },
            }
        ]

        output_file = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_file)

        result, errors = compose_pptx(
            template_files=[self.template_path],
            slide_specs=slide_specs,
            global_context=self.context,
            output=output_file,
            use_tagged_layouts=True,
        )

        # Should succeed without errors
        self.assertIsNotNone(result)
        self.assertIsNone(errors)
        self.assertTrue(os.path.exists(output_file))

        # Verify the output
        prs = Presentation(output_file)
        self.assertEqual(len(prs.slides), 1)

        # Most importantly, verify slide dimensions are within PowerPoint limits
        # Convert from EMUs to inches (914400 EMUs = 1 inch)
        width_inches = prs.slide_width / 914400
        height_inches = prs.slide_height / 914400

        # Without scaling, this graph would be ~108" x 86", which exceeds the 56" limit
        # Verify it has been scaled to fit
        self.assertLessEqual(
            width_inches,
            56,
            f"Slide width {width_inches:.2f} exceeds PowerPoint's 56 inch limit",
        )
        self.assertLessEqual(
            height_inches,
            56,
            f"Slide height {height_inches:.2f} exceeds PowerPoint's 56 inch limit",
        )

        # Verify the graph was actually scaled down (not just using minimum size)
        # The width should be close to 56 inches since it's the limiting dimension
        self.assertGreater(
            width_inches, 50, "Graph should be scaled to use available space"
        )


if __name__ == "__main__":
    unittest.main()

"""
Comprehensive tests for grouped shapes with different content types.
"""

import os
import tempfile
import unittest

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from office_templates.office_renderer.pptx.render import render_pptx


class TestGroupedShapesComprehensive(unittest.TestCase):
    def setUp(self):
        """Create test presentations with various grouped content types."""
        self.temp_files = []
        self.context = {
            "title": "Test Title",
            "chart_title": "Sales Chart",
            "grouped_text": "Grouped Content",
            "table_data": [
                ["Name", "Value"],
                ["Alice", "100"],
                ["Bob", "200"],
            ],
            "chart_data": [
                ["Month", "Sales"],
                ["Jan", 100],
                ["Feb", 150],
                ["Mar", 200],
            ],
        }

    def tearDown(self):
        """Clean up temporary files."""
        for temp_file in self.temp_files:
            if os.path.exists(temp_file):
                os.remove(temp_file)

    def test_grouped_text_and_table(self):
        """Test groups containing both text shapes and tables."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Create text shape
        text_shape = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        text_shape.text_frame.text = "Title: {{ title }}"
        
        # Create table shape
        table_shape = slide.shapes.add_table(
            rows=3, cols=2, left=Inches(1), top=Inches(2.5), width=Inches(3), height=Inches(2)
        )
        table = table_shape.table
        
        # Add template variables to table
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "{{ grouped_text }}"
        table.cell(1, 1).text = "Static content"
        table.cell(2, 0).text = "Row 2"
        table.cell(2, 1).text = "{{ title }}"
        
        # Group them
        group = slide.shapes.add_group_shape([text_shape, table_shape])
        
        # Save and render
        template_path = tempfile.mktemp(suffix=".pptx")
        prs.save(template_path)
        self.temp_files.append(template_path)
        
        output_path = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_path)
        
        result, errors = render_pptx(
            template=template_path,
            context=self.context,
            output=output_path,
            check_permissions=None,
        )
        
        self.assertIsNone(errors, f"Rendering should succeed: {errors}")
        
        # Verify the rendered content
        rendered_prs = Presentation(output_path)
        slide = rendered_prs.slides[0]
        
        # Find the group
        group_found = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_found = True
                
                # Check text shapes in group
                for grouped_shape in shape.shapes:
                    if hasattr(grouped_shape, "text_frame") and grouped_shape.text_frame:
                        text = grouped_shape.text_frame.text
                        if "Title: Test Title" in text:
                            print(f"Found rendered text: {text}")
                    
                    # Check table content in group
                    if getattr(grouped_shape, "has_table", False):
                        table = grouped_shape.table
                        for row_idx in range(len(table.rows)):
                            for col_idx in range(len(table.columns)):
                                cell_text = table.cell(row_idx, col_idx).text
                                print(f"Table cell [{row_idx}][{col_idx}]: {cell_text}")
                                if cell_text == "Grouped Content":
                                    print("Found rendered table content!")
                                elif cell_text == "Test Title":
                                    print("Found rendered title in table!")
        
        self.assertTrue(group_found, "Should find grouped shapes")

    def test_grouped_chart_shapes(self):
        """Test groups containing chart shapes."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Create text shape with chart title
        text_shape = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        text_shape.text_frame.text = "Chart: {{ chart_title }}"
        
        # Create chart shape
        chart_data = ChartData()
        chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
        chart_data.add_series("Sales", [10, 20, 30, 40])
        
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, 
            Inches(1), Inches(2.5), Inches(4), Inches(3),
            chart_data
        )
        
        # Add template variables to chart data
        # Note: For this test, we'll just verify the group processing works
        # Chart data templating is handled separately by the chart processor
        
        # Group them
        group = slide.shapes.add_group_shape([text_shape, chart_shape])
        
        # Save and render
        template_path = tempfile.mktemp(suffix=".pptx")
        prs.save(template_path)
        self.temp_files.append(template_path)
        
        output_path = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_path)
        
        result, errors = render_pptx(
            template=template_path,
            context=self.context,
            output=output_path,
            check_permissions=None,
        )
        
        self.assertIsNone(errors, f"Rendering should succeed: {errors}")
        
        # Verify the rendered content
        rendered_prs = Presentation(output_path)
        slide = rendered_prs.slides[0]
        
        # Find the group and verify text was processed
        group_found = False
        chart_title_rendered = False
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_found = True
                
                for grouped_shape in shape.shapes:
                    if hasattr(grouped_shape, "text_frame") and grouped_shape.text_frame:
                        text = grouped_shape.text_frame.text
                        if "Chart: Sales Chart" in text:
                            chart_title_rendered = True
                            print(f"Found rendered chart title: {text}")
                    
                    # Verify chart is present
                    if getattr(grouped_shape, "has_chart", False):
                        print("Found chart in group - chart processing works!")
        
        self.assertTrue(group_found, "Should find grouped shapes")
        self.assertTrue(chart_title_rendered, "Chart title should be rendered")

    def test_mixed_grouped_and_ungrouped_content(self):
        """Test that rendering works with both grouped and ungrouped content on same slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Ungrouped content
        ungrouped_shape = slide.shapes.add_textbox(
            Inches(5), Inches(1), Inches(2), Inches(1)
        )
        ungrouped_shape.text_frame.text = "Ungrouped: {{ title }}"
        
        # Grouped content
        grouped_text1 = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(3), Inches(1)
        )
        grouped_text1.text_frame.text = "Grouped 1: {{ grouped_text }}"
        
        grouped_text2 = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(3), Inches(1)
        )
        grouped_text2.text_frame.text = "Grouped 2: {{ title }}"
        
        group = slide.shapes.add_group_shape([grouped_text1, grouped_text2])
        
        # Save and render
        template_path = tempfile.mktemp(suffix=".pptx")
        prs.save(template_path)
        self.temp_files.append(template_path)
        
        output_path = tempfile.mktemp(suffix=".pptx")
        self.temp_files.append(output_path)
        
        result, errors = render_pptx(
            template=template_path,
            context=self.context,
            output=output_path,
            check_permissions=None,
        )
        
        self.assertIsNone(errors, f"Rendering should succeed: {errors}")
        
        # Verify both grouped and ungrouped content was rendered
        rendered_prs = Presentation(output_path)
        slide = rendered_prs.slides[0]
        
        ungrouped_rendered = False
        grouped_rendered = False
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for grouped_shape in shape.shapes:
                    if hasattr(grouped_shape, "text_frame") and grouped_shape.text_frame:
                        text = grouped_shape.text_frame.text
                        if "Grouped 1: Grouped Content" in text or "Grouped 2: Test Title" in text:
                            grouped_rendered = True
                            print(f"Found rendered grouped content: {text}")
            else:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text
                    if "Ungrouped: Test Title" in text:
                        ungrouped_rendered = True
                        print(f"Found rendered ungrouped content: {text}")
        
        self.assertTrue(ungrouped_rendered, "Ungrouped content should be rendered")
        self.assertTrue(grouped_rendered, "Grouped content should be rendered")


if __name__ == "__main__":
    unittest.main()
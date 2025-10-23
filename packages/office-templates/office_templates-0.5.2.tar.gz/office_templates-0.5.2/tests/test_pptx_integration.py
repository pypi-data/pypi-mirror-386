import datetime
import os
import tempfile
import unittest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from office_templates.office_renderer import render_pptx


# Dummy objects for integration testing.
class DummyUser:
    def __init__(self, name, email, is_active=True):
        self.name = name
        self.email = email
        self.is_active = is_active
        self._meta = True  # Simulate a Django model

    def __str__(self):
        return self.name


class DummyCohort:
    def __init__(self, name):
        self.name = name

    def __str__(self):
        return self.name


class DummyRequestUser:
    def has_perm(self, perm, obj):
        # Deny permission if object's name contains "deny"
        if hasattr(obj, "name") and "deny" in obj.name.lower():
            return False
        return True


class TestRendererIntegration(unittest.TestCase):
    def setUp(self):
        # Create a minimal PPTX with one text box and one table.
        self.prs = Presentation()
        blank_slide = self.prs.slide_layouts[5]
        self.slide = self.prs.slides.add_slide(blank_slide)
        # Add a text box with mixed text.
        textbox = self.slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(4), Inches(1)
        )
        textbox.text_frame.text = "Welcome, {{ user.name }}. Program: {{ program.name }}."
        self.textbox_index = len(self.slide.shapes) - 1
        # Add a table with one row, one column.
        rows, cols = 1, 1
        left, top, width, height = Inches(0.5), Inches(2), Inches(4), Inches(0.8)
        table_shape = self.slide.shapes.add_table(rows, cols, left, top, width, height)
        # Use the original mixed text with prefix.
        table_cell = table_shape.table.cell(0, 0)
        table_cell.text = "Here: {{ program.users.email }}"
        self.table_shape_index = None
        for idx, shape in enumerate(self.slide.shapes):
            if getattr(shape, "has_table", False) and shape.has_table:
                self.table_shape_index = idx
                break

        # Save this PPTX to a temporary file.
        self.temp_input = tempfile.mktemp(suffix=".pptx")
        self.temp_output = tempfile.mktemp(suffix=".pptx")
        self.prs.save(self.temp_input)

        # Set up the context.
        self.cohort = DummyCohort("Cohort A")
        self.user = DummyUser("Alice", "alice@example.com", is_active=True)
        self.user2 = DummyUser("Bob", "bob@example.com", is_active=True)
        self.user3 = DummyUser("Carol", "carol@example.com", is_active=False)
        self.program = {
            "name": "Test Program",
            "users": [self.user, self.user2, self.user3],
        }
        self.context = {
            "user": self.user,
            "program": self.program,
            "date": datetime.date(2020, 1, 15),
        }
        self.request_user = DummyRequestUser()

    def tearDown(self):
        if os.path.exists(self.temp_input):
            os.remove(self.temp_input)
        if os.path.exists(self.temp_output):
            os.remove(self.temp_output)

    def test_integration_renderer(self):
        # Run the renderer integration.
        rendered, errors = render_pptx(
            self.temp_input,
            self.context,
            self.temp_output,
            None,
        )
        self.assertIsNone(errors)

        # Open the rendered PPTX.
        prs_out = Presentation(rendered)
        # Test text box content.
        textbox = prs_out.slides[0].shapes[self.textbox_index]
        txt = textbox.text_frame.text
        self.assertIn("Welcome, Alice.", txt)
        self.assertIn("Program: Test Program", txt)
        # Check table content.
        table_shape = prs_out.slides[0].shapes[self.table_shape_index]
        self.assertIsNotNone(table_shape)
        # The original cell text is "Here: {{ program.users.email }}", so after rendering
        # it should expand into as many rows as there are program.users.
        expected_users = [u.email for u in self.program["users"]]
        # Get all rows from the table.
        rows = list(table_shape.table.rows)
        self.assertEqual(len(rows), len(expected_users))
        # For each row, check that its cell text equals "Here: " and then the corresponding email.
        for i, row in enumerate(rows):
            cell_text = row.cells[0].text.strip()
            expected = f"Here: {expected_users[i]}"
            self.assertEqual(cell_text, expected)

    def test_table_cell_formatting_preserved_on_row_expansion(self):
        # Create a PPTX with a table cell with formatting
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        table_shape = slide.shapes.add_table(
            1, 1, Inches(0.5), Inches(2), Inches(4), Inches(0.8)
        )
        cell = table_shape.table.cell(0, 0)
        cell.text = "Here: {{ program.users.email }}"
        # Set formatting: bold, font size, color
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0]
        run.font.bold = True
        run.font.size = 240000  # 24 pt
        run.font.color.rgb = RGBColor(255, 0, 0)  # Red
        # Save template
        temp_input = tempfile.mktemp(suffix=".pptx")
        temp_output = tempfile.mktemp(suffix=".pptx")
        prs.save(temp_input)
        # Context with multiple users
        program = {
            "name": "Test Program",
            "users": [
                DummyUser("Alice", "alice@example.com"),
                DummyUser("Bob", "bob@example.com"),
            ],
        }
        context = {"user": program["users"][0], "program": program}
        # Render
        rendered, errors = render_pptx(temp_input, context, temp_output, None)
        self.assertIsNone(errors)
        prs_out = Presentation(rendered)
        table_shape_out = None
        for shape in prs_out.slides[0].shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                table_shape_out = shape
                break
        self.assertIsNotNone(table_shape_out)
        rows = list(table_shape_out.table.rows)
        self.assertEqual(len(rows), 2)
        # Check formatting in each row
        for row in rows:
            cell = row.cells[0]
            para = cell.text_frame.paragraphs[0]
            run = para.runs[0]
            self.assertTrue(run.font.bold)
            # Font sizes may have small variations in rendering, allow for small difference
            self.assertAlmostEqual(run.font.size, 240000, delta=1000)
            # Color may be RGBColor, so compare as tuple
            rgb = tuple(run.font.color.rgb) if run.font.color.rgb else None
            self.assertEqual(rgb, (255, 0, 0))
        os.remove(temp_input)
        os.remove(temp_output)

    def test_callable_math_expression_integration(self):
        """Test that math operations on callable results are evaluated in PPTX rendering."""

        class DummyCallable:
            def multiply(self, a, b):
                return a * b

            def add(self, a, b):
                return a + b

        # Add a new slide with a text box containing callable math expressions
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(4), Inches(1))
        textbox.text_frame.text = (
            "Result1: {{ dummy.multiply(3,4) * 10 }}\n"
            "Result2: {{ dummy.multiply(3,4) + 5 }}\n"
            "Result3: {{ dummy.add(10,5) * 2 }}"
        )
        self.prs.save(self.temp_input)
        # Add dummy callable to context
        self.context["dummy"] = DummyCallable()
        rendered, errors = render_pptx(
            self.temp_input, self.context, self.temp_output, None
        )
        self.assertIsNone(errors)
        prs_out = Presentation(rendered)
        # The new slide is the last one
        out_textbox = prs_out.slides[-1].shapes[-1]
        txt = out_textbox.text_frame.text
        self.assertIn("Result1: 120.0", txt)  # 3*4=12, 12*10=120
        self.assertIn("Result2: 17.0", txt)  # 3*4=12, 12+5=17
        self.assertIn("Result3: 30.0", txt)  # 10+5=15, 15*2=30

    def test_image_download_and_replace(self):
        """A text box starting with %imagesqueeze% should be replaced by a picture."""

        from PIL import Image

        img = Image.new("RGB", (2, 2), color="red")
        img_file = tempfile.mktemp(suffix=".png")
        img.save(img_file)

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(2), Inches(2))
        box.text_frame.text = f"%imagesqueeze% file://{img_file}"

        self.prs.save(self.temp_input)

        rendered, errors = render_pptx(
            self.temp_input, self.context, self.temp_output, None
        )
        self.assertIsNone(errors)

        prs_out = Presentation(rendered)
        pic_shape = prs_out.slides[-1].shapes[-1]
        self.assertEqual(pic_shape.shape_type, MSO_SHAPE_TYPE.PICTURE)
        self.assertEqual(pic_shape.width, box.width)
        self.assertEqual(pic_shape.height, box.height)

        os.remove(img_file)

    def test_image_download_and_fit(self):
        """%image% should fit inside the shape keeping aspect ratio."""

        from PIL import Image

        img = Image.new("RGB", (200, 100), color="red")  # Use different dimensions to test aspect ratio
        img_file = tempfile.mktemp(suffix=".png")
        img.save(img_file)

        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(4), Inches(1))
        box.text_frame.text = f"%image% file://{img_file}"

        self.prs.save(self.temp_input)

        rendered, errors = render_pptx(
            self.temp_input, self.context, self.temp_output, None
        )
        self.assertIsNone(errors)

        prs_out = Presentation(rendered)
        pic_shape = prs_out.slides[-1].shapes[-1]
        self.assertEqual(pic_shape.shape_type, MSO_SHAPE_TYPE.PICTURE)
        self.assertEqual(pic_shape.height, box.height)
        self.assertLess(pic_shape.width, box.width)

        os.remove(img_file)


if __name__ == "__main__":
    unittest.main()
